"""
Wavenet 提案 PPT 產生器
Template-based: loads real Wavenet PPTX, keeps first 10 slides,
then appends content slides that visually continue the template style.

Template design language (reverse-engineered from real PPTX):
  - Primary navy  : #002338
  - Accent orange : #ED7D31
  - Teal          : #0097A7
  - Font          : Microsoft YaHei
  - Slide size    : 9144000 × 5143500 EMU (10 × 5.625 in)
  - Layout        : number badge top-left + title, white bg
"""
import os
import copy
import json
from io import BytesIO
from datetime import datetime

from lxml import etree as _et
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

_TEMPLATES_DIR  = os.path.join(os.path.dirname(__file__), "templates")
# The Wavenet company-intro base is ALWAYS used — never replaced by uploads.
_BASE_TEMPLATE  = os.path.join(_TEMPLATES_DIR, "wavenet_template.pptx")
# Design tokens extracted from reference uploads.
_DESIGN_TOKENS  = os.path.join(_TEMPLATES_DIR, "design_tokens.json")


def _load_design_tokens() -> dict:
    """Return color/font tokens extracted from the reference file, or {}."""
    if not os.path.exists(_DESIGN_TOKENS):
        return {}
    try:
        return json.load(open(_DESIGN_TOKENS, encoding="utf-8"))
    except Exception:
        return {}


def extract_design_tokens(pptx_path: str) -> dict:
    """
    Scan content slides (11+) of a reference PPTX for dominant colors/font.
    Returns { "primary", "accent", "secondary", "font", "source" }.
    """
    from collections import Counter
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return {}

    color_counts: Counter = Counter()
    font_counts:  Counter = Counter()
    start = min(10, len(prs.slides))

    for slide in prs.slides[start:]:
        for shape in slide.shapes:
            try:
                rgb = shape.fill.fore_color.rgb
                r, g, b = rgb[0], rgb[1], rgb[2]
                # Skip near-white and near-black
                if not (r > 230 and g > 230 and b > 230) and \
                   not (r < 30  and g < 30  and b < 30):
                    color_counts[f"#{r:02X}{g:02X}{b:02X}"] += 1
            except Exception:
                pass
            if hasattr(shape, "text_frame"):
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                font_counts[run.font.name] += 1
                except Exception:
                    pass

    top_colors = [c for c, _ in color_counts.most_common(5)]
    top_font   = font_counts.most_common(1)[0][0] if font_counts else None
    tokens: dict = {}
    if len(top_colors) > 0: tokens["primary"]   = top_colors[0]
    if len(top_colors) > 1: tokens["accent"]    = top_colors[1]
    if len(top_colors) > 2: tokens["secondary"] = top_colors[2]
    if top_font:             tokens["font"]      = top_font
    tokens["source"] = os.path.basename(pptx_path)
    return tokens

# ── Slide geometry matching actual Wavenet template (10 × 5.625 in) ──────────
SW   = Emu(9144000)    # 10 in
SH   = Emu(5143500)    # 5.625 in
ML   = Inches(0.3)     # left margin
MR   = Inches(0.3)     # right margin
UW   = SW - ML - MR   # usable width = 9.4 in

# ── Wavenet brand palette ─────────────────────────────────────────────────────
NAVY    = RGBColor(0x00, 0x23, 0x38)
ORANGE  = RGBColor(0xED, 0x7D, 0x31)
TEAL    = RGBColor(0x00, 0x97, 0xA7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x21, 0x21, 0x21)
GRAY    = RGBColor(0x75, 0x75, 0x75)
L_NAVY  = RGBColor(0xE3, 0xE8, 0xEC)
L_ORNG  = RGBColor(0xFD, 0xEF, 0xE4)
L_TEAL  = RGBColor(0xE0, 0xF5, 0xF7)
L_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MIDGRAY = RGBColor(0xE0, 0xE0, 0xE0)

FONT = "Microsoft YaHei"

# ── Key y-positions (all proportioned for 5.625" height) ─────────────────────
CONT_T = Inches(0.78)                      # content area top (below header+underline)
CARD_B = SH - Inches(0.35)                 # max card bottom (above footer)


# ─────────────────────────── Low-level helpers ────────────────────────────────

def _rect(slide, l, t, w, h, color: RGBColor):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _box(slide, l, t, w, h, *,
         fill: RGBColor = None,
         text: str = "",
         size: int = 10,
         bold: bool = False,
         color: RGBColor = None,
         align: PP_ALIGN = PP_ALIGN.LEFT,
         ml: float = 0.1,
         mt: float = 0.05):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        shp.fill.background()
        shp.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(ml)
    tf.margin_right  = Inches(ml)
    tf.margin_top    = Inches(mt)
    tf.margin_bottom = Inches(mt)
    tf._txBody.bodyPr.set('anchor', 't')
    if text:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or DARK
        r.font.name = FONT
    return shp, tf


def _para(tf, text: str, *,
          size: int = 9,
          bold: bool = False,
          color: RGBColor = None,
          align: PP_ALIGN = PP_ALIGN.LEFT,
          space_before: int = 0,
          italic: bool = False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if text:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color or DARK
        r.font.name = FONT
    return p


def _trunc(s, n):
    if not s:
        return ""
    s = str(s)
    return s[:n] + ("..." if len(s) > n else "")


# ─────────────────────────── Shared design chrome ─────────────────────────────

def _page_header(slide, number: str, title: str,
                 number_bg: RGBColor = NAVY,
                 title_color: RGBColor = NAVY):
    badge_w = badge_h = Inches(0.42)
    badge_y = Inches(0.15)
    _rect(slide, ML, badge_y, badge_w, badge_h, number_bg)
    _box(slide, ML, badge_y, badge_w, badge_h,
         text=number, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, ml=0, mt=0.04)
    _box(slide, ML + badge_w + Inches(0.1), badge_y,
         UW - badge_w - Inches(0.1), badge_h,
         text=title, size=15, bold=True, color=title_color, ml=0.05, mt=0.05)
    _rect(slide, ML, Inches(0.62), UW, Inches(0.025), ORANGE)


def _footer(slide):
    _box(slide, ML, SH - Inches(0.22), UW, Inches(0.20),
         text="潮網科技 Wavenet Technology  ·  Digital Marketing Proposal",
         size=7, color=GRAY, align=PP_ALIGN.RIGHT, ml=0, mt=0.03)


def _section_label(slide, text: str, l, t, color: RGBColor = NAVY):
    _box(slide, l, t, UW, Inches(0.18),
         text=text.upper(), size=7, bold=True, color=color, ml=0, mt=0)


# ─────────────────────────── Slide helpers ────────────────────────────────────

def _delete_slides_from(prs: Presentation, keep_count: int):
    sldIdLst = prs.slides._sldIdLst
    NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    while len(prs.slides) > keep_count:
        idx = len(prs.slides) - 1
        rId = sldIdLst[idx].get(NS)
        if rId:
            prs.part.drop_rel(rId)
        del sldIdLst[idx]


def _add_slide(prs: Presentation, ref_slide=None):
    if ref_slide is not None:
        return _copy_slide(prs, ref_slide)
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[0])


def _copy_slide(dest_prs: Presentation, src_slide):
    """
    Copy visual-only (non-text) elements from src_slide into a new slide in dest_prs.
    Shapes that contain text are skipped so the caller can layer fresh AI content on top.
    Image relationships are remapped where possible.
    """
    layout = dest_prs.slide_layouts[-1]
    for lay in dest_prs.slide_layouts:
        if "blank" in lay.name.lower():
            layout = lay
            break
    new_slide = dest_prs.slides.add_slide(layout)

    # Copy slide background fill
    try:
        src_bg = src_slide.background._element
        dst_bg = new_slide.background._element
        for child in list(dst_bg):
            dst_bg.remove(child)
        for child in list(src_bg):
            dst_bg.append(copy.deepcopy(child))
    except Exception:
        pass

    # Copy only shapes that have no text content (decorative/structural shapes)
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    src_tree = src_slide.shapes._spTree
    dst_tree = new_slide.shapes._spTree
    while len(dst_tree) > 2:
        dst_tree.remove(dst_tree[-1])

    for elem in list(src_tree)[2:]:
        has_text = any(
            (t.text or "").strip()
            for t in elem.iter(f"{{{NS_A}}}t")
        )
        if has_text:
            continue  # caller adds fresh content; skip original text shapes
        dst_tree.append(copy.deepcopy(elem))

    # Remap image/media relationship IDs
    try:
        rId_map = {}
        for rId, rel in src_slide.part.rels.items():
            try:
                if rel.is_external:
                    new_rId = new_slide.part.relate_to(
                        rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                if new_rId != rId:
                    rId_map[rId] = new_rId
            except Exception:
                pass
        if rId_map:
            xml_str = _et.tostring(dst_tree).decode("utf-8")
            for old_id, new_id in rId_map.items():
                xml_str = xml_str.replace(f'"{old_id}"', f'"{new_id}"')
            new_tree = _et.fromstring(xml_str.encode("utf-8"))
            while len(dst_tree):
                dst_tree.remove(dst_tree[-1])
            for child in new_tree:
                dst_tree.append(child)
    except Exception:
        pass

    return new_slide


# ─────────────────────────── Cover update ─────────────────────────────────────

def _update_cover(slide, proposal: dict):
    company  = proposal.get("company_name", "")
    title    = proposal.get("title", "數位行銷策略提案")
    date_str = datetime.now().strftime("%Y%m%d")
    subs = {
        "富悅國際行銷": company,
        "有限公司": "",
        "20260316": date_str,
        "數位行銷規劃": title,
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in subs.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


# ─────────────────────────── Section divider ──────────────────────────────────

def _slide_section(prs, section_num: str, title: str, accent: RGBColor = NAVY):
    slide = _add_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _rect(slide, Inches(0), Inches(0), Inches(0.14), SH, accent)
    _box(slide, Inches(0.24), Inches(0.75), Inches(1.0), Inches(0.9),
         text=section_num, size=52, bold=True, color=accent, ml=0, mt=0)
    _box(slide, Inches(0.24), Inches(1.75), UW - Inches(0.24), Inches(0.55),
         text=title, size=22, bold=True, color=DARK, ml=0, mt=0)
    _rect(slide, Inches(0.24), Inches(2.36), Inches(1.5), Inches(0.045), accent)
    _box(slide, ML, SH - Inches(0.28), UW, Inches(0.24),
         text="潮網科技 Wavenet Technology",
         size=8, color=GRAY, align=PP_ALIGN.RIGHT, ml=0, mt=0)


# ─────────────────────────── Phase 2 slides ───────────────────────────────────

def _slide_phase2a(prs, p2: dict, phase_num: str = "01", ref_slide=None):
    slide = _add_slide(prs, ref_slide=ref_slide)
    if ref_slide is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

    _page_header(slide, phase_num, "全漏斗策略規劃 — 現況診斷與方向")
    _footer(slide)

    diag     = _trunc(p2.get("current_diagnosis", ""), 120)
    approach = _trunc(p2.get("recommended_approach", ""), 120)
    insight  = _trunc(p2.get("key_insight", ""), 80)

    half_w = (UW - Inches(0.15)) / 2
    ty = CONT_T

    _rect(slide, ML, ty, half_w, Inches(0.26), NAVY)
    _box(slide, ML, ty, half_w, Inches(0.26),
         text="客戶現況診斷", size=9, bold=True, color=WHITE, ml=0.1, mt=0.05)
    box_h = CARD_B - (ty + Inches(0.26)) - (Inches(0.45) if insight else 0)
    _, tf = _box(slide, ML, ty + Inches(0.26), half_w, box_h, fill=L_NAVY)
    _para(tf, diag, size=9, color=DARK)

    rx = ML + half_w + Inches(0.15)
    _rect(slide, rx, ty, half_w, Inches(0.26), ORANGE)
    _box(slide, rx, ty, half_w, Inches(0.26),
         text="策略方向建議", size=9, bold=True, color=WHITE, ml=0.1, mt=0.05)
    _, tf = _box(slide, rx, ty + Inches(0.26), half_w, box_h, fill=L_ORNG)
    _para(tf, approach, size=9, color=DARK)

    if insight:
        bar_t = CARD_B - Inches(0.38)
        _rect(slide, ML, bar_t, Inches(0.05), Inches(0.38), ORANGE)
        _box(slide, ML + Inches(0.1), bar_t, UW - Inches(0.1), Inches(0.38),
             fill=L_ORNG,
             text=f"Key Insight  |  {insight}",
             size=9, bold=True, color=NAVY, ml=0.1, mt=0.08)


def _slide_phase2b(prs, p2: dict, phase_num: str = "02", ref_slide=None):
    slide = _add_slide(prs, ref_slide=ref_slide)
    if ref_slide is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

    _page_header(slide, phase_num, "四階段全漏斗媒體規劃")
    _footer(slide)

    funnel = p2.get("funnel", [])
    palette = [
        (NAVY,   L_NAVY),
        (TEAL,   L_TEAL),
        (ORANGE, L_ORNG),
        (RGBColor(0x3C, 0x5A, 0xA0), RGBColor(0xE8, 0xED, 0xF7)),
    ]

    n_cols  = 4
    gap     = Inches(0.1)
    col_w   = (UW - gap * (n_cols - 1)) / n_cols
    hdr_h   = Inches(0.28)
    cont_t  = CONT_T + hdr_h
    cont_h  = CARD_B - cont_t

    for i, stage in enumerate(funnel[:4]):
        lx = ML + i * (col_w + gap)
        fg, bg = palette[i % 4]

        _rect(slide, lx, CONT_T, col_w, hdr_h, fg)
        _box(slide, lx, CONT_T, col_w, hdr_h,
             text=_trunc(stage.get("stage", f"Phase {i+1}"), 8),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, ml=0.04, mt=0.04)

        _, tf = _box(slide, lx, cont_t, col_w, cont_h, fill=bg, ml=0.07, mt=0.07)
        _para(tf, "目標", size=7, bold=True, color=fg)
        _para(tf, _trunc(stage.get("objective", ""), 50), size=8, color=DARK)
        _para(tf, "媒體管道", size=7, bold=True, color=fg, space_before=3)
        channels = "、".join(stage.get("channels", []))
        _para(tf, _trunc(channels, 38), size=8, color=DARK)
        _para(tf, "目標受眾", size=7, bold=True, color=fg, space_before=3)
        _para(tf, _trunc(stage.get("audience", ""), 38), size=8, color=DARK)
        _para(tf, "KPI", size=7, bold=True, color=fg, space_before=3)
        _para(tf, _trunc(stage.get("kpi", ""), 32), size=8, bold=True, color=fg)


# ─────────────────────────── Phase 3 ──────────────────────────────────────────

def _slide_phase3(prs, p3: dict, phase_num: str = "03", ref_slide=None):
    slide = _add_slide(prs, ref_slide=ref_slide)
    if ref_slide is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

    _page_header(slide, phase_num, "市場數據洞察", TEAL, TEAL)
    _footer(slide)

    bm = p3.get("benchmarks", {})
    kpi_items = [
        ("ROAS", bm.get("industry_avg_roas", "—")),
        ("CPC",  bm.get("industry_avg_cpc",  "—")),
        ("CTR",  bm.get("industry_avg_ctr",  "—")),
        ("CVR",  bm.get("industry_avg_cvr",  "—")),
    ]

    n_cols = 4
    gap    = Inches(0.1)
    kpi_w  = (UW - gap * (n_cols - 1)) / n_cols
    kpi_h  = Inches(0.82)

    for i, (label, val) in enumerate(kpi_items):
        lx = ML + i * (kpi_w + gap)
        _rect(slide, lx, CONT_T, kpi_w, Inches(0.06), TEAL)
        _, tf = _box(slide, lx, CONT_T + Inches(0.06), kpi_w, kpi_h, fill=L_TEAL, ml=0.05, mt=0.06)
        _para(tf, _trunc(str(val), 12), size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        _para(tf, f"產業平均 {label}", size=7, color=GRAY, align=PP_ALIGN.CENTER, space_before=2)

    ty2    = CONT_T + kpi_h + Inches(0.16)
    half_w = (UW - Inches(0.14)) / 2
    box_h  = CARD_B - (ty2 + Inches(0.2))

    gap_text = _trunc(p3.get("competitive_gap", ""), 100)
    _section_label(slide, "競爭差距分析", ML, ty2, TEAL)
    _, tf = _box(slide, ML, ty2 + Inches(0.2), half_w, box_h, fill=L_TEAL, ml=0.08, mt=0.08)
    _para(tf, gap_text, size=8, color=DARK)

    opps = p3.get("growth_opportunities", [])
    rx = ML + half_w + Inches(0.14)
    _section_label(slide, "成長機會", rx, ty2, ORANGE)
    _, tf = _box(slide, rx, ty2 + Inches(0.2), half_w, box_h, fill=L_ORNG, ml=0.08, mt=0.08)
    for opp in opps[:4]:
        _para(tf, f"- {_trunc(opp, 40)}", size=8, color=DARK)


# ─────────────────────────── Phase 4 ──────────────────────────────────────────

def _slide_phase4(prs, p4: dict, phase_num: str = "04", ref_slide=None):
    slide = _add_slide(prs, ref_slide=ref_slide)
    if ref_slide is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

    PURPLE = RGBColor(0x5C, 0x35, 0xA8)
    L_PURP = RGBColor(0xEF, 0xEC, 0xF9)
    _page_header(slide, phase_num, "廣告創意策略", PURPLE, PURPLE)
    _footer(slide)

    formats = p4.get("recommended_formats", [])
    fx    = ML
    tag_h = Inches(0.26)
    _section_label(slide, "建議廣告格式", ML, CONT_T, PURPLE)
    for fmt in formats[:5]:
        tag_w = Inches(1.3)
        if fx + tag_w > ML + UW:
            break
        _rect(slide, fx, CONT_T + Inches(0.2), tag_w, tag_h, L_PURP)
        _box(slide, fx, CONT_T + Inches(0.2), tag_w, tag_h,
             text=_trunc(fmt, 12), size=8, bold=True,
             color=PURPLE, align=PP_ALIGN.CENTER, ml=0.04, mt=0.03)
        fx += tag_w + Inches(0.1)

    dims = p4.get("creative_dimensions", {})
    dim_data = [
        ("trust_building", "品牌信任", NAVY,   L_NAVY),
        ("pain_point",     "痛點訴求", ORANGE, L_ORNG),
        ("conversion",     "轉換促購", TEAL,   L_TEAL),
    ]
    n_cols = 3
    gap    = Inches(0.1)
    col_w  = (UW - gap * (n_cols - 1)) / n_cols
    ty2    = CONT_T + Inches(0.55)
    hdr_h  = Inches(0.28)
    body_h = CARD_B - (ty2 + hdr_h)

    for i, (key, default_name, fg, bg) in enumerate(dim_data):
        lx  = ML + i * (col_w + gap)
        dim = dims.get(key, {})
        name = _trunc(dim.get("name", default_name), 8)
        desc = _trunc(dim.get("description", ""), 80)
        examples = dim.get("examples", [])

        _rect(slide, lx, ty2, col_w, hdr_h, fg)
        _box(slide, lx, ty2, col_w, hdr_h,
             text=name, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, ml=0.04, mt=0.05)

        _, tf = _box(slide, lx, ty2 + hdr_h, col_w, body_h, fill=bg, ml=0.08, mt=0.07)
        _para(tf, desc, size=8, color=DARK)
        if examples:
            _para(tf, "素材範例", size=7, bold=True, color=fg, space_before=3)
            for ex in examples[:2]:
                _para(tf, f"- {_trunc(ex, 30)}", size=7, color=GRAY)


# ─────────────────────────── Phase 5 slides ───────────────────────────────────

def _slide_phase5a(prs, p5: dict, phase_num: str = "05", ref_slide=None):
    slide = _add_slide(prs, ref_slide=ref_slide)
    if ref_slide is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

    _page_header(slide, phase_num, "媒體預算配置", ORANGE, ORANGE)
    _footer(slide)

    budget = _trunc(str(p5.get("monthly_budget", "—")), 14)
    roas   = _trunc(str(p5.get("expected_roas",  "—")), 14)
    alloc  = p5.get("channel_allocation", [])

    pill_w = Inches(1.65)
    pill_h = Inches(0.7)
    for j, (val, label, c) in enumerate([
        (budget, "月預算規模", NAVY),
        (roas,   "預期 ROAS", TEAL),
    ]):
        px = ML + j * (pill_w + Inches(0.15))
        _rect(slide, px, CONT_T, pill_w, pill_h, c)
        _box(slide, px, CONT_T, pill_w, Inches(0.42),
             text=val, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, ml=0.04, mt=0.06)
        _box(slide, px, CONT_T + Inches(0.42), pill_w, Inches(0.28),
             text=label, size=7, color=RGBColor(0xCC, 0xDD, 0xE5),
             align=PP_ALIGN.CENTER, ml=0.04, mt=0.02)

    ty       = CONT_T + pill_h + Inches(0.1)
    _section_label(slide, "媒體管道配置", ML, ty, ORANGE)
    ty      += Inches(0.2)

    bar_max  = UW * 0.4
    row_h    = Inches(0.38)
    ch_col_w = Inches(1.1)
    max_rows = int((CARD_B - ty) / row_h)

    for ch in alloc[:min(6, max_rows)]:
        channel = _trunc(ch.get("channel", ""), 10)
        pct     = min(int(ch.get("percentage", 0)), 100)
        reason  = _trunc(ch.get("rationale", ""), 30)

        _box(slide, ML, ty, ch_col_w, row_h,
             text=channel, size=8, bold=True, color=DARK, ml=0, mt=0.08)

        bar_x = ML + ch_col_w + Inches(0.08)
        _rect(slide, bar_x, ty + Inches(0.09), bar_max, Inches(0.2), MIDGRAY)
        if pct > 0:
            _rect(slide, bar_x, ty + Inches(0.09),
                  bar_max * pct / 100, Inches(0.2), ORANGE)
        _box(slide, bar_x + bar_max + Inches(0.07), ty, Inches(0.42), row_h,
             text=f"{pct}%", size=8, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER, ml=0, mt=0.08)
        remaining_w = UW - ch_col_w - Inches(0.08) - bar_max - Inches(0.55)
        if remaining_w > Inches(0.3):
            _box(slide, bar_x + bar_max + Inches(0.55), ty, remaining_w, row_h,
                 text=reason, size=7, color=GRAY, ml=0, mt=0.08)
        ty += row_h


def _slide_phase5b(prs, p5: dict, phase_num: str = "06", ref_slide=None):
    slide = _add_slide(prs, ref_slide=ref_slide)
    if ref_slide is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

    _page_header(slide, phase_num, "關鍵節點與執行時程", ORANGE, ORANGE)
    _footer(slide)

    campaigns = p5.get("key_campaigns", [])
    tl_note   = _trunc(p5.get("timeline_note", ""), 80)

    ty  = CONT_T
    cws = [Inches(1.8), Inches(1.3), UW - Inches(3.3)]
    hdrs = ["活動名稱", "執行時間", "重點說明"]
    lx  = ML
    for h, w in zip(hdrs, cws):
        _rect(slide, lx, ty, w - Inches(0.03), Inches(0.28), NAVY)
        _box(slide, lx, ty, w - Inches(0.03), Inches(0.28),
             text=h, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             ml=0.05, mt=0.04)
        lx += w

    ty    += Inches(0.28)
    row_h  = Inches(0.37)
    max_rows = int((CARD_B - ty - Inches(0.42)) / row_h)

    for i, c in enumerate(campaigns[:min(5, max_rows)]):
        bg = L_GRAY if i % 2 == 0 else WHITE
        lx = ML
        vals = [_trunc(c.get("name",  ""), 16),
                _trunc(c.get("timing",""), 12),
                _trunc(c.get("focus", ""), 45)]
        for val, w in zip(vals, cws):
            _rect(slide, lx, ty, w - Inches(0.03), row_h, bg)
            _box(slide, lx, ty, w - Inches(0.03), row_h,
                 text=val, size=8, color=DARK, ml=0.06, mt=0.07)
            lx += w
        ty += row_h

    if tl_note:
        note_t = min(ty + Inches(0.08), CARD_B - Inches(0.35))
        _rect(slide, ML, note_t, UW, Inches(0.32), L_ORNG)
        _box(slide, ML + Inches(0.08), note_t, UW - Inches(0.08), Inches(0.32),
             text=f"執行建議  |  {tl_note}", size=8, color=NAVY, ml=0.06, mt=0.07)
        ty = note_t + Inches(0.38)

    cta_t = SH - Inches(0.62)
    _rect(slide, Inches(0), cta_t, SW, Inches(0.52), NAVY)
    _box(slide, ML, cta_t, UW, Inches(0.52),
         text="立即聯繫潮網科技  ·  讓我們為您打造專屬數位行銷方案",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, ml=0.08, mt=0.12)


# ─────────────────────────── Main entry ───────────────────────────────────────

def generate_pptx(proposal: dict) -> BytesIO:
    # ── Apply design tokens from reference file (if any) ──────────────────────
    # Always use the Wavenet base template — uploads only supply colour/font.
    global NAVY, ORANGE, TEAL, FONT
    orig_navy, orig_orange, orig_teal, orig_font = NAVY, ORANGE, TEAL, FONT

    def _hex(s: str) -> RGBColor:
        h = s.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    tokens = _load_design_tokens()
    if tokens.get("primary"):   NAVY   = _hex(tokens["primary"])
    if tokens.get("accent"):    ORANGE = _hex(tokens["accent"])
    if tokens.get("secondary"): TEAL   = _hex(tokens["secondary"])
    if tokens.get("font"):      FONT   = tokens["font"]

    # Load reference content slides (11+) from the activated design reference file
    ref_slides = []
    if tokens.get("source"):
        ref_path = os.path.join(_TEMPLATES_DIR, tokens["source"])
        if os.path.exists(ref_path):
            try:
                ref_prs = Presentation(ref_path)
                ref_slides = list(ref_prs.slides)[10:]  # slides 11+ are content slides
            except Exception:
                pass

    def _ref(i):
        return ref_slides[i % len(ref_slides)] if ref_slides else None

    try:
        prs = Presentation(_BASE_TEMPLATE)
        _update_cover(prs.slides[0], proposal)
        _delete_slides_from(prs, keep_count=10)

        content = proposal.get("content") or {}
        p2 = content.get("phase2") or {}
        p3 = content.get("phase3") or {}
        p4 = content.get("phase4") or {}
        p5 = content.get("phase5") or {}

        _slide_section(prs, "01", "數位行銷策略提案", NAVY)
        _slide_phase2a(prs, p2, "01", ref_slide=_ref(0))
        _slide_phase2b(prs, p2, "02", ref_slide=_ref(1))
        _slide_section(prs, "02", "市場數據洞察", TEAL)
        _slide_phase3(prs, p3, "03", ref_slide=_ref(2))
        _slide_section(prs, "03", "廣告創意策略", RGBColor(0x5C, 0x35, 0xA8))
        _slide_phase4(prs, p4, "04", ref_slide=_ref(3))
        _slide_section(prs, "04", "媒體預算規劃", ORANGE)
        _slide_phase5a(prs, p5, "05", ref_slide=_ref(4))
        _slide_phase5b(prs, p5, "06", ref_slide=_ref(5))

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf
    finally:
        # Restore palette so module state is clean for the next request
        NAVY, ORANGE, TEAL, FONT = orig_navy, orig_orange, orig_teal, orig_font
