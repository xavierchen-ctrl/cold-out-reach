import os
import io
import json
import re
from datetime import timedelta
from typing import Optional
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from fastapi.responses import StreamingResponse
from pptx import Presentation as PptxPresentation
from pydantic import BaseModel
from sqlalchemy.orm import Session
import httpx
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from openai import OpenAI
from database import get_db
from models import User, Lead, LeadActivity, LeadStatus, ActivityType
from schemas import DraftRequest, DraftResponse
from auth import get_current_user
from utils import now_tw

router = APIRouter(prefix="/api/ai", tags=["ai"])

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")


def _openai_chat(prompt: str, model: str = "gpt-4o-mini") -> str:
    oai = OpenAI(api_key=OPENAI_API_KEY)
    return oai.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
    ).choices[0].message.content.strip()


def _openai_json(prompt: str, model: str = "gpt-4o-mini") -> dict:
    oai = OpenAI(api_key=OPENAI_API_KEY)
    return json.loads(oai.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    ).choices[0].message.content)


TEMPLATES = {
    "intro": "初次開發信。這是第一次接觸，請保持專業但友善，簡要介紹我們的數位行銷服務，並提出一個明確的行動呼籲（安排 15 分鐘通話）。",
    "followup": "追蹤跟進信。我們之前已有接觸但尚未回應，請溫和提醒，不要施壓，強調我們能提供的具體價值。",
    "proposal": "報價提案信。針對客戶需求提出具體方案，包含服務項目概述、預期效益，並邀請進一步討論細節。",
}


@router.post("/draft", response_model=DraftResponse)
def generate_draft(
    body: DraftRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    template_hint = TEMPLATES.get(body.template_type, TEMPLATES["intro"])

    background_section = ""
    if body.customer_background and body.customer_background.strip():
        background_section = f"""
客戶官網／背景資訊（請仔細閱讀並融入信件，讓內容更貼近客戶實際情況）：
---
{body.customer_background.strip()[:3000]}
---
"""

    prompt = f"""你是一位專業的數位行銷業務，正在撰寫開發信。

目標客戶資訊：
- 公司名稱：{lead.company_name}
- 聯絡人：{lead.contact_name or "（未知）"}
- 職稱：{lead.title or "（未知）"}
- 產業：{lead.industry or "數位行銷"}
- 城市：{lead.city or "台灣"}
- 官網：{lead.website or "（未知）"}
{background_section}
信件類型：{body.template_type}
指導方針：{template_hint}

請用繁體中文撰寫一封專業的開發信。若有提供客戶背景資訊，請具體提及客戶的產品/服務，讓信件更有針對性，而非泛泛而談。
回應格式（嚴格遵守）：
SUBJECT: <信件主旨>
BODY:
<信件內文>

不要加任何其他說明。"""

    try:
        text = _openai_chat(prompt)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI error: {str(e)}")

    subject = ""
    body_text = ""
    if "SUBJECT:" in text and "BODY:" in text:
        parts = text.split("BODY:", 1)
        subject = parts[0].replace("SUBJECT:", "").strip()
        body_text = parts[1].strip()
    else:
        lines = text.split("\n")
        subject = lines[0]
        body_text = "\n".join(lines[1:]).strip()

    return DraftResponse(subject=subject, body=body_text)


# ── AI Enrich ─────────────────────────────────────────────────────────────────

class EnrichRequest(BaseModel):
    company_name: str
    lead_id: str


import re as _re
from urllib.parse import urljoin as _urljoin

_EMAIL_RE = _re.compile(r'[\w.+\-]+@[\w\-]+\.[a-zA-Z]{2,}')

_PHONE_RE = _re.compile(
    r'(?:'
    r'\(?0[2-8]\)?[-\s\.]{0,2}\d{3,4}[-\s\.]{0,2}\d{4}'
    r'|\(?09\d{2}\)?[-\s\.]{0,2}\d{3}[-\s\.]{0,2}\d{3}'
    r'|\+886[-\s\.]{0,2}[2-9][-\s\.]{0,2}\d{3,4}[-\s\.]{0,2}\d{4}'
    r')'
)

_CONTACT_KWS = [
    "聯絡我們", "contact-us", "contactus", "contact_us",
    "聯絡", "連絡", "contact",
    "reach", "get-in-touch", "connect",
    "about", "關於",
]

_CONTACT_PATHS = [
    "/contact", "/contact-us", "/contactus", "/contact_us",
    "/about/contact", "/about-us", "/about",
    "/pages/contact", "/pages/contact-us",
    "/zh/contact", "/tw/contact",
    "/get-in-touch", "/support",
    "/company/contact",
]

_SKIP_EMAIL_PARTS = ["example", "test", "noreply", "no-reply", "donotreply",
                     ".png", ".jpg", ".gif", ".svg", ".css", ".js", "sentry",
                     "wixpress", "w3school", "schema.org"]

_HTTP_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
        "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120 Safari/537.36"
    )
}


def _strip_tags(html: str) -> str:
    text = _re.sub(r'<[^>]+>', ' ', html)
    return _re.sub(r'\s+', ' ', text).strip()


def _find_emails(text: str) -> list[str]:
    normalized = (
        text
        .replace(" [at] ", "@").replace(" (at) ", "@")
        .replace("[at]", "@").replace("(at)", "@")
        .replace(" [dot] ", ".").replace(" (dot) ", ".")
        .replace("[dot]", ".").replace("(dot)", ".")
    )
    found = _EMAIL_RE.findall(normalized)
    result = []
    seen: set[str] = set()
    for e in found:
        el = e.lower()
        if el in seen:
            continue
        if any(skip in el for skip in _SKIP_EMAIL_PARTS):
            continue
        seen.add(el)
        result.append(e)
    return result


def _find_phones(text: str) -> list[str]:
    found = _PHONE_RE.findall(text)
    cleaned = []
    seen: set[str] = set()
    for p in found:
        c = _re.sub(r'[\s\-\.\(\)]', '', p)
        if c not in seen and len(c) >= 8:
            seen.add(c)
            cleaned.append(p.strip())
    return cleaned


def _extract_jsonld_contacts(html: str) -> dict:
    schema_re = _re.compile(
        r'<script[^>]+type=["\']application/ld\+json["\'][^>]*>(.*?)</script>',
        _re.DOTALL | _re.IGNORECASE,
    )
    emails: list[str] = []
    phones: list[str] = []
    for match in schema_re.findall(html):
        try:
            data = json.loads(match.strip())
            items = data if isinstance(data, list) else [data]
            for item in items:
                for key in ("email", "Email"):
                    val = item.get(key)
                    if val and isinstance(val, str):
                        emails.append(val)
                for key in ("telephone", "Telephone", "phone", "Phone"):
                    val = item.get(key)
                    if val and isinstance(val, str):
                        phones.append(val)
                cp = item.get("contactPoint") or []
                if isinstance(cp, dict):
                    cp = [cp]
                for point in cp:
                    if point.get("email"):
                        emails.append(point["email"])
                    if point.get("telephone"):
                        phones.append(point["telephone"])
        except Exception:
            pass
    seen_e: set[str] = set()
    seen_p: set[str] = set()
    emails = [e for e in emails if not (e.lower() in seen_e or seen_e.add(e.lower()))]  # type: ignore[func-returns-value]
    phones = [p for p in phones if not (_re.sub(r'[^\d]', '', p) in seen_p or seen_p.add(_re.sub(r'[^\d]', '', p)))]  # type: ignore[func-returns-value]
    return {"emails": emails, "phones": phones}


async def _scrape_contact_info(website: str) -> dict:
    base = website if website.startswith("http") else f"https://{website}"

    try:
        async with httpx.AsyncClient(
            timeout=15, headers=_HTTP_HEADERS, follow_redirects=True
        ) as client:
            try:
                res = await client.get(base)
                main_html = res.text
            except Exception:
                return {}

            jsonld = _extract_jsonld_contacts(main_html)
            if jsonld["emails"] or jsonld["phones"]:
                return {
                    "email": jsonld["emails"][0] if jsonld["emails"] else None,
                    "phone": jsonld["phones"][0] if jsonld["phones"] else None,
                    "contact_url": None,
                }

            link_re = _re.compile(r'href=["\']([^"\'\s#][^"\']*)["\']', _re.IGNORECASE)
            all_links = link_re.findall(main_html)

            contact_url: str | None = None
            best_score = 0
            for href in all_links:
                href_lower = href.lower()
                if href_lower.startswith(("mailto:", "tel:", "javascript:")):
                    continue
                score = 0
                for i, kw in enumerate(_CONTACT_KWS):
                    if kw.lower() in href_lower:
                        score += max(5 - i, 1)
                if score > best_score:
                    best_score = score
                    contact_url = _urljoin(base, href)

            if not contact_url:
                for suffix in _CONTACT_PATHS:
                    try:
                        r = await client.get(_urljoin(base, suffix), timeout=8)
                        if r.status_code == 200 and len(r.text) > 200:
                            contact_url = _urljoin(base, suffix)
                            break
                    except Exception:
                        pass

            target_html = main_html
            if contact_url:
                try:
                    r = await client.get(contact_url, timeout=10)
                    if r.status_code == 200:
                        target_html = r.text
                        jsonld = _extract_jsonld_contacts(r.text)
                        if jsonld["emails"] or jsonld["phones"]:
                            return {
                                "email": jsonld["emails"][0] if jsonld["emails"] else None,
                                "phone": jsonld["phones"][0] if jsonld["phones"] else None,
                                "contact_url": contact_url,
                            }
                except Exception:
                    pass

            text = _strip_tags(target_html)
            emails = _find_emails(text)
            phones = _find_phones(text)

            # ── Step 7: OpenAI fallback if regex finds nothing ────────────
            if not emails and not phones and OPENAI_API_KEY:
                snippet = text[:4000]
                contact_prompt = f"""以下是一個公司聯絡頁面的文字內容，請仔細尋找 Email 地址和電話號碼（包含台灣常見格式如 (02)2345-6789 或 0912-345-678）。
只回傳 JSON，格式如下（沒有則填 null）：
{{"email": "第一個有效 email", "phone": "第一個有效電話"}}

內容：
{snippet}"""
                try:
                    parsed = _openai_json(contact_prompt)
                    if parsed.get("email"):
                        emails = [parsed["email"]]
                    if parsed.get("phone"):
                        phones = [parsed["phone"]]
                except Exception:
                    pass

            return {
                "email": emails[0] if emails else None,
                "phone": phones[0] if phones else None,
                "contact_url": contact_url,
            }
    except Exception:
        return {}


@router.post("/enrich")
async def enrich_company(
    body: EnrichRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """Use OpenAI to enrich company info (industry, city, company_size, email, phone)."""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    prompt = f"""請根據公司名稱「{body.company_name}」，推測並補全以下資訊（繁體中文，JSON格式回傳）：

{{
  "industry": "產業類別（例如：科技/製造/零售/金融/餐飲/教育等）",
  "city": "總部城市（台灣城市名稱）",
  "company_size": "公司規模（例如：1-10人/11-50人/51-200人/201-500人/500人以上）",
  "website": "公司官方網站完整網址（若能確定請填，否則填 null）",
  "summary": "公司簡介（50字以內）"
}}

只回傳JSON，不要任何說明。"""

    try:
        data = _openai_json(prompt)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="AI 回傳格式錯誤")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI error: {str(e)}")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if lead:
        if data.get("industry") and not lead.industry:
            lead.industry = data["industry"]
        if data.get("city") and not lead.city:
            lead.city = data["city"]
        if data.get("company_size") and not lead.company_size:
            lead.company_size = data["company_size"]

        ai_website = data.get("website")
        if ai_website and not lead.website:
            if isinstance(ai_website, str) and "." in ai_website:
                if not ai_website.startswith("http"):
                    ai_website = f"https://{ai_website}"
                lead.website = ai_website
                data["ai_suggested_website"] = ai_website

        website_to_scrape = lead.website
        contact_info: dict = {}
        if website_to_scrape:
            contact_info = await _scrape_contact_info(website_to_scrape)
            scraped_email = contact_info.get("email")
            scraped_phone = contact_info.get("phone")
            if scraped_email and not lead.email:
                lead.email = scraped_email
                data["scraped_email"] = scraped_email
            if scraped_phone and not lead.phone:
                lead.phone = scraped_phone
                data["scraped_phone"] = scraped_phone
            data["contact_url"] = contact_info.get("contact_url")
        else:
            data["scraped_email"] = None
            data["scraped_phone"] = None
            data["contact_url"] = None

        db.commit()

    return data


# ── Pipeline Health ───────────────────────────────────────────────────────────

@router.post("/pipeline_health")
def pipeline_health(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """Analyze overall pipeline and return Markdown health report."""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    from sqlalchemy import func

    total = db.query(Lead).count()
    by_status = {}
    for row in db.query(Lead.status, func.count()).group_by(Lead.status).all():
        by_status[row[0].value] = row[1]

    week_ago = now_tw() - timedelta(days=7)
    emails_week = db.query(LeadActivity).filter(
        LeadActivity.type == ActivityType.email_sent,
        LeadActivity.created_at >= week_ago,
    ).count()

    stale_threshold = now_tw() - timedelta(days=7)
    stale_count = 0
    for lead in db.query(Lead).filter(Lead.status.in_([LeadStatus.contacted, LeadStatus.replied])).all():
        last = db.query(LeadActivity).filter(
            LeadActivity.lead_id == lead.id
        ).order_by(LeadActivity.created_at.desc()).first()
        if not last or last.created_at < stale_threshold:
            stale_count += 1

    prompt = f"""你是一位資深 B2B 業務顧問。請根據以下 Pipeline 數據，生成一份繁體中文的 Markdown 健診報告：

## Pipeline 數據
- 總名單數：{total}
- 各狀態分佈：{json.dumps(by_status, ensure_ascii=False)}
- 本週發信數：{emails_week}
- 超過7天未跟進：{stale_count} 筆

## 報告要求
請用 Markdown 格式，包含：
1. **健康度評分**（0-100分 + 評語）
2. **關鍵風險**（條列）
3. **建議行動**（優先順序排列）
4. **本週重點任務**（3項）

保持專業、具體、可執行。字數約 300-500 字。"""

    try:
        return {"report": _openai_chat(prompt)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI error: {str(e)}")


# ── Gmail Reply Detection ─────────────────────────────────────────────────────

@router.post("/gmail/check_replies")
def check_gmail_replies(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """Check Gmail inbox for replies and auto-update lead status."""
    if not current_user.gmail_token:
        raise HTTPException(status_code=400, detail="Gmail not connected")

    import base64
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build

    token_data = json.loads(current_user.gmail_token)
    creds = Credentials(
        token=token_data["token"],
        refresh_token=token_data.get("refresh_token"),
        token_uri=token_data.get("token_uri", "https://oauth2.googleapis.com/token"),
        client_id=token_data.get("client_id"),
        client_secret=token_data.get("client_secret"),
        scopes=token_data.get("scopes"),
    )

    try:
        service = build("gmail", "v1", credentials=creds)
        results = service.users().messages().list(
            userId="me", q="in:inbox is:unread", maxResults=50
        ).execute()
        messages = results.get("messages", [])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gmail error: {str(e)}")

    leads = db.query(Lead).filter(
        Lead.assigned_to == current_user.id,
        Lead.email.isnot(None),
        Lead.status.in_([LeadStatus.contacted]),
    ).all()
    lead_emails = {l.email.lower(): l for l in leads if l.email}

    updated = []
    for msg in messages:
        try:
            full_msg = service.users().messages().get(
                userId="me", id=msg["id"], format="metadata",
                metadataHeaders=["From", "Subject"]
            ).execute()
            headers = {h["name"]: h["value"] for h in full_msg.get("payload", {}).get("headers", [])}
            from_email = headers.get("From", "").lower()

            for email, lead in lead_emails.items():
                if email in from_email:
                    if lead.status == LeadStatus.contacted:
                        lead.status = LeadStatus.replied
                        activity = LeadActivity(
                            lead_id=lead.id,
                            type=ActivityType.status_change,
                            content=f"Gmail 自動偵測回信 → 狀態更新為「已回覆」\nSubject: {headers.get('Subject', '')}",
                            created_by=current_user.id,
                        )
                        db.add(activity)
                        updated.append(lead.company_name)
                        break
        except Exception:
            continue

    if updated:
        db.commit()

    return {"updated_leads": updated, "checked_messages": len(messages)}


# ── AI 客製化 Email（輸入官網生成） ───────────────────────────────────────────

class GenerateEmailRequest(BaseModel):
    website_url: str
    product: str
    lead_id: Optional[str] = None
    tone: str = "professional"


PRODUCT_DESCRIPTIONS = {
    "SEO優化": "透過技術SEO、內容優化與外部連結建立，提升網站搜尋排名，增加自然流量",
    "廣告投放": "專業管理Google/Meta/LINE等數位廣告，精準觸達目標受眾，最大化廣告投資報酬率",
    "社群代操": "代管Facebook、Instagram、LINE官帳號等社群媒體，建立品牌形象與粉絲互動",
    "整合行銷": "跨平台整合行銷策略，結合SEO、社群、廣告、EDM等多管道協同運作",
    "KOL行銷": "媒合適合品牌的網紅KOL，執行口碑行銷活動，擴大品牌聲量與轉換",
}

TONE_DESCRIPTIONS = {
    "professional": "專業正式，商務語氣，展現專業能力",
    "friendly": "親切友善，拉近距離，自然對話風格",
    "urgent": "製造急迫感，限時優惠或市場競爭壓力",
}

_STRIP_TAGS_RE = re.compile(r'<[^>]+>')


def strip_html_tags(html: str) -> str:
    text = _STRIP_TAGS_RE.sub(' ', html)
    return re.sub(r'\s+', ' ', text).strip()


@router.post("/generate-email")
async def generate_email(
    body: GenerateEmailRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """Generate a customized email using the client's website content."""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    website_summary = ""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 Chrome/120 Safari/537.36",
        }
        async with httpx.AsyncClient(timeout=15, headers=headers, follow_redirects=True) as client:
            res = await client.get(body.website_url)
            raw_html = res.text[:8000]
            website_summary = strip_html_tags(raw_html)[:3000]
    except Exception as e:
        website_summary = f"（無法讀取官網：{str(e)[:100]}）"

    product_desc = PRODUCT_DESCRIPTIONS.get(body.product, body.product)
    tone_desc = TONE_DESCRIPTIONS.get(body.tone, body.tone)

    prompt = f"""你是潮網科技的業務，正在撰寫一封客製化開發信。

【客戶官網內容摘要】
{website_summary[:2000]}

【潮網主推服務】
服務名稱：{body.product}
服務說明：{product_desc}

【語氣要求】
{tone_desc}

請根據以上資訊，用繁體中文撰寫一封 200 字以內的開發信。
信件必須：
1. 提及客戶的業務特色（從官網摘要中找出1-2個具體亮點）
2. 說明潮網如何幫助客戶成長
3. 提出明確的下一步（如：安排 15 分鐘線上會議）
4. 語氣符合要求

回應格式（嚴格遵守）：
SUBJECT: <信件主旨>
BODY:
<信件內文>

不要加任何其他說明。"""

    try:
        text = _openai_chat(prompt)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI error: {str(e)}")

    subject = ""
    body_text = ""
    if "SUBJECT:" in text and "BODY:" in text:
        parts = text.split("BODY:", 1)
        subject = parts[0].replace("SUBJECT:", "").strip()
        body_text = parts[1].strip()
    else:
        lines = text.split("\n")
        subject = lines[0]
        body_text = "\n".join(lines[1:]).strip()

    return {
        "subject": subject,
        "body": body_text,
        "website_summary": website_summary[:500],
    }


# ── 提案信生成 ─────────────────────────────────────────────────────────────────

class GenerateProposalRequest(BaseModel):
    lead_id: str
    product: str
    tone: str = "professional"
    include_benchmark: bool = True


@router.post("/generate-proposal")
async def generate_proposal(
    body: GenerateProposalRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    score = lead.enriched_score or 0

    has_gtm = tech.get("gtm", False)
    has_pixel = tech.get("meta_pixel", False)
    has_ga4 = tech.get("ga4", False)
    has_meta_ads = ad.get("meta", {}).get("has_ads", False)
    has_google_ads = ad.get("google_ads", {}).get("has_ads", False)

    PRODUCTS = {
        "廣告投放": "Meta/Google 廣告投放優化，提升 ROAS",
        "SEO優化": "搜尋引擎優化，提升自然流量",
        "社群代操": "FB/IG/LINE 社群經營，提升品牌聲量",
        "整合行銷": "跨通路整合行銷方案，串聯線上線下",
        "KOL行銷": "網紅/KOL 合作行銷，精準觸及目標受眾",
        "程序化廣告": "Programmatic 廣告投放，精準鎖定受眾",
    }
    product_desc = PRODUCTS.get(body.product, body.product)
    tone_map = {"professional": "專業正式", "friendly": "親切友善", "urgent": "急迫有力"}
    tone_label = tone_map.get(body.tone, "專業正式")

    prompt = f"""
你是潮網科技（Wavenet Technology）的資深業務顧問。
請為以下潛在客戶撰寫一封專業的開發提案信（繁體中文）。

## 客戶資訊
- 公司：{lead.company_name}
- 聯絡人：{lead.contact_name or '您好'}
- 職稱：{lead.title or '行銷主管'}
- 產業：{lead.industry or '數位行銷'}
- 官網：{lead.website or '未知'}
- 含金量分數：{score}/100

## 客戶現況分析（從官網偵測）
- Google Tag Manager：{"已安裝" if has_gtm else "未安裝"}
- Meta Pixel：{"已安裝" if has_pixel else "未安裝"}
- GA4：{"已安裝" if has_ga4 else "未安裝"}
- Meta 廣告：{"有在投放" if has_meta_ads else "未偵測到"}
- Google 廣告：{"有在投放" if has_google_ads else "未偵測到"}

## 主推服務
{body.product}：{product_desc}

## 提案信要求
1. 主旨（subject）：吸引注意，提及客戶公司名稱，不超過 15 字
2. 正文（body）：
   - 開場：提及對該公司的了解（基於現況分析）
   - 現況痛點：指出 1-2 個可改善的地方（根據偵測結果）
   - 解決方案：說明潮網如何幫助（主推 {body.product}）
   - 成功案例：簡短提及潮網過往成效（可虛構合理數字）
   - CTA：約定時間電話/視訊討論，留下聯絡方式
   - 簽名：潮網科技業務團隊

語氣：{tone_label}
長度：300-400 字

請輸出 JSON：
{{"subject": "信件主旨", "body": "信件正文（保留換行）", "key_points": ["重點1", "重點2", "重點3"]}}
"""

    try:
        data = _openai_json(prompt)
        return {
            "subject": data.get("subject", ""),
            "body": data.get("body", ""),
            "key_points": data.get("key_points", []),
            "lead_id": body.lead_id,
            "product": body.product,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 生成失敗：{str(e)}")


# ── PPT 簡報背景資料生成 ───────────────────────────────────────────────────────

class PptBriefRequest(BaseModel):
    lead_id: str


@router.post("/ppt-brief")
async def generate_ppt_brief(
    body: PptBriefRequest,
    db: Session = Depends(get_db),
    _current_user: User = Depends(get_current_user),
):
    """根據廠商資料生成簡報所需的背景資料（Markdown 格式，可複製到 Gamma/Canva 等工具）。"""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    score = lead.enriched_score or 0

    tech_parts = []
    if tech.get("ga4"): tech_parts.append("GA4 分析")
    if tech.get("gtm"): tech_parts.append("Google Tag Manager")
    if tech.get("meta_pixel"): tech_parts.append("Meta Pixel")
    if tech.get("google_ads_tag"): tech_parts.append("Google Ads 追蹤")
    tech_summary = "、".join(tech_parts) if tech_parts else "尚未偵測到追蹤工具"

    ad_parts = []
    if ad.get("meta", {}).get("has_ads"): ad_parts.append("Meta 廣告")
    if ad.get("google_ads", {}).get("has_ads"): ad_parts.append("Google 廣告")
    ad_summary = "、".join(ad_parts) if ad_parts else "尚未偵測到廣告投放"

    prompt = f"""你是潮網科技（Wavenet Technology）的資深業務顧問，正在準備一份拜訪客戶前的簡報背景資料。

請根據以下廠商資訊，用繁體中文生成一份適合製作 PowerPoint 的結構化背景資料。
每個章節請用清晰的標題和條列式重點，語氣專業，便於直接貼入 Gamma.app / Canva 等 AI 簡報工具。

## 廠商資訊
- 公司名稱：{lead.company_name}
- 聯絡人：{lead.contact_name or "未知"}
- 職稱：{lead.title or "未知"}
- 產業：{lead.industry or "未知"}
- 城市：{lead.city or "台灣"}
- 公司規模：{lead.company_size or "未知"}
- 官方網站：{lead.website or "未知"}
- 含金量評分：{score}/100
- 已安裝追蹤工具：{tech_summary}
- 廣告投放現況：{ad_summary}
- 備注：{lead.notes or "無"}

## 輸出格式（請嚴格依照此架構）

# 【{lead.company_name}】簡報背景資料

## 1. 公司概況
（3-5 個條列，涵蓋：產業定位、主要業務、服務對象、地理範圍）

## 2. 推測主要產品／服務
（根據產業與公司名稱推測，3-5 個條列）

## 3. 數位行銷現況分析
（根據偵測到的工具與廣告投放分析，3-5 個條列，點出現況強項與缺口）

## 4. 市場機會與痛點
（3-4 個條列，說明該廠商可能面臨的行銷挑戰）

## 5. 潮網建議合作方向
（3-4 個條列，具體列出可切入的服務項目與預期效益）

## 6. 簡報開場白建議
（1 段 50 字以內的開場白，適合面對面拜訪時使用）

---
請直接輸出 Markdown 內容，不要加任何前言或說明。"""

    try:
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
        )
        return {"brief": response.choices[0].message.content.strip(), "company_name": lead.company_name}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 生成失敗：{str(e)}")


# ── Google Slides 簡報建立 ─────────────────────────────────────────────────────

@router.post("/create-slides")
async def create_google_slides(
    body: PptBriefRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """用 OpenAI 生成簡報內容，並透過 Google Slides API 在用戶 Drive 建立簡報。"""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    if not current_user.gmail_token:
        raise HTTPException(status_code=400, detail="請先在設定頁面連結 Google 帳號")

    token_data = json.loads(current_user.gmail_token)
    if "https://www.googleapis.com/auth/presentations" not in token_data.get("scopes", []):
        raise HTTPException(
            status_code=403,
            detail="需要重新連結 Google 帳號以授予 Google Slides 存取權限，請前往設定頁面重新連結"
        )

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    # 1. Gemini 生成結構化簡報內容
    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    score = lead.enriched_score or 0

    tech_parts = []
    if tech.get("ga4"): tech_parts.append("GA4")
    if tech.get("gtm"): tech_parts.append("GTM")
    if tech.get("meta_pixel"): tech_parts.append("Meta Pixel")
    tech_str = "、".join(tech_parts) if tech_parts else "無"

    ad_parts = []
    if ad.get("meta", {}).get("has_ads"): ad_parts.append("Meta 廣告")
    if ad.get("google_ads", {}).get("has_ads"): ad_parts.append("Google 廣告")
    ad_str = "、".join(ad_parts) if ad_parts else "無"

    prompt = f"""你是潮網科技的業務顧問。請根據以下廠商資訊，生成一份 6 頁 Google Slides 簡報的內容。

廠商資訊：
- 公司名稱：{lead.company_name}
- 產業：{lead.industry or "未知"}
- 城市：{lead.city or "台灣"}
- 官網：{lead.website or "未知"}
- 公司規模：{lead.company_size or "未知"}
- 含金量評分：{score}/100
- 已安裝追蹤工具：{tech_str}
- 廣告投放：{ad_str}

只輸出 JSON，不要其他說明：
{{
  "title": "【{lead.company_name}】數位行銷提案",
  "slides": [
    {{"title": "關於 {lead.company_name}", "bullets": ["要點1", "要點2", "要點3"]}},
    {{"title": "主要產品與服務", "bullets": ["項目1", "項目2", "項目3"]}},
    {{"title": "數位行銷現況分析", "bullets": ["分析1", "分析2", "分析3"]}},
    {{"title": "市場機會與挑戰", "bullets": ["機會1", "機會2", "機會3"]}},
    {{"title": "潮網建議合作方向", "bullets": ["建議1", "建議2", "建議3"]}},
    {{"title": "下一步行動", "bullets": ["安排 15 分鐘線上會議", "提供詳細提案與報價", "聯絡信箱：service@wavenet.com.tw"]}}
  ]
}}

每個 bullets 包含 3-5 個繁體中文條列，具體且有說服力。"""

    try:
        from openai import OpenAI
        oai = OpenAI(api_key=OPENAI_API_KEY)
        response = oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
        )
        raw = response.choices[0].message.content.strip()
        slide_content = json.loads(raw)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 內容生成失敗：{str(e)}")

    # 2. 建立 Google Slides 簡報
    try:
        creds = Credentials(
            token=token_data["token"],
            refresh_token=token_data.get("refresh_token"),
            token_uri=token_data.get("token_uri", "https://oauth2.googleapis.com/token"),
            client_id=token_data.get("client_id"),
            client_secret=token_data.get("client_secret"),
            scopes=token_data.get("scopes"),
        )

        slides_service = build("slides", "v1", credentials=creds)

        presentation = slides_service.presentations().create(
            body={"title": slide_content["title"]}
        ).execute()
        pres_id = presentation["presentationId"]
        default_slide_id = presentation["slides"][0]["objectId"]

        requests = [{"deleteObject": {"objectId": default_slide_id}}]

        for i, slide_data in enumerate(slide_content.get("slides", [])):
            slide_id = f"slide_{i}"
            title_id = f"title_{i}"
            body_id = f"body_{i}"

            requests.append({
                "createSlide": {
                    "objectId": slide_id,
                    "insertionIndex": i,
                    "slideLayoutReference": {"predefinedLayout": "BLANK"},
                }
            })

            requests += [
                {
                    "createShape": {
                        "objectId": title_id,
                        "shapeType": "TEXT_BOX",
                        "elementProperties": {
                            "pageObjectId": slide_id,
                            "size": {
                                "height": {"magnitude": 900000, "unit": "EMU"},
                                "width": {"magnitude": 8229600, "unit": "EMU"},
                            },
                            "transform": {
                                "scaleX": 1, "scaleY": 1,
                                "translateX": 457200, "translateY": 300000,
                                "unit": "EMU",
                            },
                        },
                    }
                },
                {
                    "insertText": {
                        "objectId": title_id,
                        "text": slide_data.get("title", ""),
                        "insertionIndex": 0,
                    }
                },
                {
                    "updateTextStyle": {
                        "objectId": title_id,
                        "style": {
                            "bold": True,
                            "fontSize": {"magnitude": 28, "unit": "PT"},
                            "foregroundColor": {
                                "opaqueColor": {
                                    "rgbColor": {"red": 0.11, "green": 0.22, "blue": 0.47}
                                }
                            },
                        },
                        "textRange": {"type": "ALL"},
                        "fields": "bold,fontSize,foregroundColor",
                    }
                },
            ]

            bullets = slide_data.get("bullets", [])
            if bullets:
                requests += [
                    {
                        "createShape": {
                            "objectId": body_id,
                            "shapeType": "TEXT_BOX",
                            "elementProperties": {
                                "pageObjectId": slide_id,
                                "size": {
                                    "height": {"magnitude": 3500000, "unit": "EMU"},
                                    "width": {"magnitude": 8229600, "unit": "EMU"},
                                },
                                "transform": {
                                    "scaleX": 1, "scaleY": 1,
                                    "translateX": 457200, "translateY": 1300000,
                                    "unit": "EMU",
                                },
                            },
                        }
                    },
                    {
                        "insertText": {
                            "objectId": body_id,
                            "text": "\n".join(f"• {b}" for b in bullets),
                            "insertionIndex": 0,
                        }
                    },
                    {
                        "updateTextStyle": {
                            "objectId": body_id,
                            "style": {"fontSize": {"magnitude": 18, "unit": "PT"}},
                            "textRange": {"type": "ALL"},
                            "fields": "fontSize",
                        }
                    },
                ]

        slides_service.presentations().batchUpdate(
            presentationId=pres_id,
            body={"requests": requests},
        ).execute()

        return {
            "url": f"https://docs.google.com/presentation/d/{pres_id}/edit",
            "title": slide_content["title"],
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Google Slides 建立失敗：{str(e)}")


# ── PPTX 模板上傳與生成 ────────────────────────────────────────────────────────

PPTX_TEMPLATES_DIR = os.getenv("TEMPLATES_DIR", "/tmp")


def _collect_text_shapes(shapes):
    """Recursively collect all shapes with text frames, including inside group shapes."""
    result = []
    for shape in shapes:
        try:
            if hasattr(shape, "shapes"):  # group shape
                result.extend(_collect_text_shapes(shape.shapes))
            elif shape.has_text_frame:
                result.append(shape)
        except Exception:
            pass
    return result


def _write_lines_to_tf(tf, lines: list):
    """Write lines into a text frame using direct XML manipulation."""
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy as _cp

    txBody = tf._txBody
    paras = txBody.findall(qn("a:p"))
    if not paras:
        return False

    first_p = paras[0]

    for i, line in enumerate(lines):
        if i < len(paras):
            p = paras[i]
        else:
            p = _cp.deepcopy(first_p)
            for r in p.findall(qn("a:r")):
                p.remove(r)
            txBody.append(p)

        runs = p.findall(qn("a:r"))
        if runs:
            t_el = runs[0].find(qn("a:t"))
            if t_el is not None:
                t_el.text = line
            for r in runs[1:]:
                p.remove(r)
        else:
            r_el = etree.SubElement(p, qn("a:r"))
            t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = line

    # Remove surplus paragraphs
    current = txBody.findall(qn("a:p"))
    for p in current[max(1, len(lines)):]:
        txBody.remove(p)
    return True


def _detect_template_company(prs) -> str:
    """Find the most frequently repeated short text across slides — likely the template company name."""
    from collections import Counter
    counter: Counter = Counter()
    for slide in prs.slides:
        slide_texts: set = set()
        for shape in _collect_text_shapes(slide.shapes):
            try:
                t = shape.text_frame.text.strip()
                if 2 <= len(t) <= 20:
                    slide_texts.add(t)
            except Exception:
                pass
        for t in slide_texts:
            counter[t] += 1
    n = len(prs.slides)
    candidates = [(t, c) for t, c in counter.items() if c >= max(2, n // 2)]
    if not candidates:
        candidates = counter.most_common(5)
    if not candidates:
        return ""
    candidates.sort(key=lambda x: (-x[1], len(x[0])))
    return candidates[0][0]


def _global_replace_text(prs, old_text: str, new_text: str) -> int:
    """Replace every occurrence of old_text in all text runs across all slides."""
    if not old_text or old_text == new_text:
        return 0
    from pptx.oxml.ns import qn
    count = 0
    for slide in prs.slides:
        for shape in _collect_text_shapes(slide.shapes):
            try:
                txBody = shape.text_frame._txBody
                for t_el in txBody.iter(qn("a:t")):
                    if t_el.text and old_text in t_el.text:
                        t_el.text = t_el.text.replace(old_text, new_text)
                        count += 1
            except Exception:
                pass
    return count


def _fill_pptx_slide(slide, title_text: str, bullets: list):
    all_shapes = _collect_text_shapes(slide.shapes)

    title_shape, body_shape = None, None
    for shape in all_shapes:
        try:
            ph = shape.placeholder_format
            if ph is None:
                continue
            if ph.idx == 0 and title_shape is None:
                title_shape = shape
            elif ph.idx in (1, 2) and body_shape is None:
                body_shape = shape
        except Exception:
            pass

    if title_shape is None or body_shape is None:
        non_ph = [s for s in all_shapes if not _has_placeholder(s)]
        # 1st try: explicit markers in template
        for s in non_ph:
            try:
                txt = s.text_frame.text
                if "{{title}}" in txt and title_shape is None:
                    title_shape = s
                if "{{content}}" in txt and body_shape is None:
                    body_shape = s
            except Exception:
                pass
        # 2nd try: largest shapes that have existing text content (skip empty bg boxes)
        if title_shape is None or body_shape is None:
            def _area(s):
                try:
                    return s.width * s.height
                except Exception:
                    return 0
            # Prefer non-empty shapes; fall back to any shape if all are empty
            non_empty = [s for s in non_ph if s.text_frame.text.strip()]
            pool = non_empty if non_empty else non_ph
            by_area = sorted(pool, key=_area, reverse=True)
            print(f"[PPTX] pool={len(pool)} (non-empty={len(non_empty)}), top areas: {[_area(s) for s in by_area[:5]]}")
            content_candidates = [s for s in by_area if _area(s) > 0]
            if len(content_candidates) >= 2:
                body_shape = content_candidates[0]
                title_shape = content_candidates[1]
            elif content_candidates:
                title_shape = content_candidates[0]

    if title_shape:
        try:
            print(f"[PPTX] title_shape text: {title_shape.text_frame.text[:40]!r}")
            _write_lines_to_tf(title_shape.text_frame, [title_text])
        except Exception as e:
            print(f"[PPTX] title fill error: {e}")

    if body_shape:
        try:
            print(f"[PPTX] body_shape text: {body_shape.text_frame.text[:40]!r}")
            _write_lines_to_tf(body_shape.text_frame, bullets)
        except Exception as e:
            print(f"[PPTX] body fill error: {e}")


def _has_placeholder(shape):
    try:
        return shape.placeholder_format is not None
    except Exception:
        return False


@router.post("/upload-pptx-template")
async def upload_pptx_template(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
):
    """上傳 PPTX 設計模板，供後續生成簡報使用。"""
    if not (file.filename or "").lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="請上傳 .pptx 格式的檔案")

    os.makedirs(PPTX_TEMPLATES_DIR, exist_ok=True)
    template_path = os.path.join(PPTX_TEMPLATES_DIR, f"pptx_tpl_{current_user.id}.pptx")

    content = await file.read()
    with open(template_path, "wb") as f:
        f.write(content)

    prs = PptxPresentation(io.BytesIO(content))
    template_company = _detect_template_company(prs)
    print(f"[PPTX] detected template_company: {template_company!r}")

    meta_path = os.path.join(PPTX_TEMPLATES_DIR, f"pptx_tpl_{current_user.id}_meta.json")
    with open(meta_path, "w", encoding="utf-8") as mf:
        json.dump({"template_company": template_company}, mf, ensure_ascii=False)

    return {
        "ok": True,
        "slide_count": len(prs.slides),
        "message": f"模板上傳成功（共 {len(prs.slides)} 頁）",
        "template_company": template_company,
    }


class GeneratePptxRequest(BaseModel):
    lead_id: str
    extra_context: str = ""
    use_template: bool = False
    context_images: list[str] = []
    client_type: str = "b2c"  # b2c | b2b | b2b_biotech


def _build_pptx_from_scratch(slides_data: list, company_name: str) -> io.BytesIO:
    """Generate a clean, professional PPTX without any template."""
    from pptx import Presentation as _Prs
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    DARK_BLUE = RGBColor(0x1A, 0x37, 0x6C)
    ACCENT    = RGBColor(0x00, 0x8C, 0xD7)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT      = RGBColor(0x1A, 0x1A, 0x2E)

    prs = _Prs()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)

        # White background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE

        # Dark blue header bar
        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = DARK_BLUE
        hdr.line.fill.background()

        # Accent stripe on left edge
        bar = slide.shapes.add_shape(1, Inches(0), Inches(1.8), Inches(0.12), Inches(5.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # Slide number badge (top-right)
        badge = slide.shapes.add_textbox(Inches(12.5), Inches(0.3), Inches(0.7), Inches(0.5))
        bp = badge.text_frame.paragraphs[0]
        bp.add_run().text = str(idx + 1)
        bp.runs[0].font.size  = Pt(13)
        bp.runs[0].font.color.rgb = RGBColor(0x88, 0xAA, 0xDD)

        # Title in header
        tbox = slide.shapes.add_textbox(Inches(0.45), Inches(0.32), Inches(12.0), Inches(1.25))
        tbox.text_frame.word_wrap = True
        tp = tbox.text_frame.paragraphs[0]
        run = tp.add_run()
        run.text = sd.get("title", "")
        run.font.size  = Pt(30)
        run.font.bold  = True
        run.font.color.rgb = WHITE

        # Company name in header (small, right-aligned below title)
        cbox = slide.shapes.add_textbox(Inches(0.45), Inches(1.42), Inches(11), Inches(0.35))
        cp = cbox.text_frame.paragraphs[0]
        crun = cp.add_run()
        crun.text = company_name
        crun.font.size  = Pt(11)
        crun.font.color.rgb = RGBColor(0x99, 0xBB, 0xEE)

        # Bullet points
        bbox = slide.shapes.add_textbox(Inches(0.45), Inches(2.1), Inches(12.7), Inches(5.1))
        bbox.text_frame.word_wrap = True
        for j, bullet in enumerate(sd.get("bullets", [])):
            p = bbox.text_frame.paragraphs[0] if j == 0 else bbox.text_frame.add_paragraph()
            r = p.add_run()
            r.text = f"▸  {bullet}"
            r.font.size  = Pt(19)
            r.font.color.rgb = TEXT
            p.space_before = Pt(10)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


@router.post("/generate-pptx")
async def generate_pptx(
    body: GeneratePptxRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """用 OpenAI 生成簡報內容，套入 PPTX 模板後回傳可下載檔案。"""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    template_path = os.path.join(PPTX_TEMPLATES_DIR, f"pptx_tpl_{current_user.id}.pptx")
    use_template = body.use_template and os.path.exists(template_path)

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    n_slides = 6  # default for scratch builds
    template_company = ""

    if use_template:
        prs_check = PptxPresentation(template_path)
        n_slides = len(prs_check.slides)
        meta_path = os.path.join(PPTX_TEMPLATES_DIR, f"pptx_tpl_{current_user.id}_meta.json")
        if os.path.exists(meta_path):
            try:
                with open(meta_path, encoding="utf-8") as mf:
                    template_company = json.load(mf).get("template_company", "")
            except Exception:
                pass
        if not template_company:
            template_company = _detect_template_company(prs_check)
        print(f"[PPTX] template_company='{template_company}'  target='{lead.company_name}'")

    # Fetch website content for accurate company description
    website_text = ""
    if lead.website:
        try:
            async with httpx.AsyncClient(timeout=10, follow_redirects=True, headers=_HTTP_HEADERS) as client:
                res = await client.get(lead.website if lead.website.startswith("http") else f"https://{lead.website}")
                website_text = _strip_tags(res.text)[:3000]
        except Exception:
            pass

    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    score = lead.enriched_score or 0
    tech_str = "、".join(k for k, v in [("GA4", tech.get("ga4")), ("GTM", tech.get("gtm")), ("Meta Pixel", tech.get("meta_pixel"))] if v) or "無"
    ad_str = "、".join(k for k, v in [("Meta 廣告", ad.get("meta", {}).get("has_ads")), ("Google 廣告", ad.get("google_ads", {}).get("has_ads"))] if v) or "無"

    website_section = ""
    if website_text:
        website_section = f"""
【官網內容摘要（請以此為準，不得自行捏造公司業務）】
{website_text}
"""

    extra_section = ""
    if body.extra_context and body.extra_context.strip():
        extra_section = f"""
【補充說明（優先參考）】
{body.extra_context.strip()}
"""

    all_topics = [
        f"關於 {lead.company_name}（公司概況、產業定位）",
        "主要產品與服務",
        "數位行銷現況分析",
        "市場機會與挑戰",
        "潮網建議合作方向",
        "下一步行動計畫",
    ]
    topics_str = "\n".join(f"{i+1}. {t}" for i, t in enumerate(all_topics[:n_slides]))

    prompt = f"""你是潮網科技的業務顧問，正在為客戶製作提案簡報。

【重要原則】
- 所有內容必須根據下方提供的資料撰寫，不得憑空捏造
- 若官網摘要有提供，以官網內容為主要依據
- 不要假設公司的產品或服務，必須從資料中找依據
{extra_section}{website_section}
【廠商基本資料】
- 公司：{lead.company_name}
- 產業：{lead.industry or "（請從官網內容判斷）"}
- 城市：{lead.city or "台灣"}
- 規模：{lead.company_size or "未知"}
- 含金量：{score}/100
- 數位追蹤工具：{tech_str}
- 廣告投放：{ad_str}

【簡報頁面主題（依序生成 {n_slides} 頁）】
{topics_str}

只輸出 JSON，不要其他說明：
{{"slides": [{{"title": "標題", "bullets": ["要點1", "要點2", "要點3"]}}]}}

每頁 3-5 個繁體中文條列，內容必須與該公司實際業務相符。"""

    try:
        from openai import OpenAI
        oai = OpenAI(api_key=OPENAI_API_KEY)
        resp = oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
        )
        slide_content = json.loads(resp.choices[0].message.content.strip())
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 內容生成失敗：{str(e)}")

    try:
        slides_data = slide_content.get("slides", [])
        from urllib.parse import quote
        safe_name = re.sub(r"[^\w一-鿿]", "_", lead.company_name)
        filename = f"{safe_name}_提案簡報.pptx"
        encoded_filename = quote(filename)

        if use_template:
            prs = PptxPresentation(template_path)

            # Step 1: fill title + body shapes per slide with AI-generated content
            for i, slide in enumerate(prs.slides):
                if i >= len(slides_data):
                    break
                _fill_pptx_slide(slide, slides_data[i].get("title", ""), slides_data[i].get("bullets", []))

            # Step 2: replace every remaining occurrence of template company name with actual company
            if template_company and template_company != lead.company_name:
                replaced = _global_replace_text(prs, template_company, lead.company_name)
                print(f"[PPTX] global replaced '{template_company}' → '{lead.company_name}': {replaced} occurrences")

            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
        else:
            print(f"[PPTX] no template — building from scratch for '{lead.company_name}'")
            output = _build_pptx_from_scratch(slides_data, lead.company_name)

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=\"proposal.pptx\"; filename*=UTF-8''{encoded_filename}"},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX 生成失敗：{str(e)}")


# ── PPTX 內容生成（Version B：前端 PptxGenJS 渲染） ────────────────────────────

@router.post("/generate-pptx-content")
async def generate_pptx_content(
    body: GeneratePptxRequest,
    db: Session = Depends(get_db),
    _current_user: User = Depends(get_current_user),
):
    """Generate AI slide content as JSON only; client renders to PPTX via PptxGenJS."""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    # Fetch website content
    website_text = ""
    if lead.website:
        try:
            async with httpx.AsyncClient(timeout=10, follow_redirects=True, headers=_HTTP_HEADERS) as client:
                url = lead.website if lead.website.startswith("http") else f"https://{lead.website}"
                res = await client.get(url)
                website_text = _strip_tags(res.text)[:3000]
        except Exception:
            pass

    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    score = lead.enriched_score or 0
    tech_str = "、".join(k for k, v in [("GA4", tech.get("ga4")), ("GTM", tech.get("gtm")), ("Meta Pixel", tech.get("meta_pixel"))] if v) or "無"
    ad_str = "、".join(k for k, v in [("Meta 廣告", ad.get("meta", {}).get("has_ads")), ("Google 廣告", ad.get("google_ads", {}).get("has_ads"))] if v) or "無"

    website_section = f"\n【官網內容摘要（請以此為準，不得自行捏造公司業務）】\n{website_text}\n" if website_text else ""
    extra_section = f"\n【補充說明（優先參考）】\n{body.extra_context.strip()}\n" if body.extra_context and body.extra_context.strip() else ""

    if body.client_type == "b2b_biotech":
        topics = [
            f"關於 {lead.company_name}（技術平台、核心競爭力）",
            "主要技術產品與服務特色",
            "目標市場與潛在合作夥伴分析",
            "市場痛點診斷",
            "B2B 購買旅程與決策者分析",
            "數位行銷現況分析",
            "年度行銷目標 KPI 設定",
            "整合媒體策略規劃",
            "預算配置方案",
            "LinkedIn ABM 決策者廣告策略",
            "Google / SEO 技術關鍵字策略",
            "思想領導力內容行銷計畫",
            "會展行銷與Geo-fencing策略",
            "季度執行時程計畫",
            "下一步行動建議",
        ]
    else:
        topics = [
            f"關於 {lead.company_name}（公司概況、品牌定位）",
            "主要產品與服務特色",
            "目標客群與市場分析",
            "品牌現況問題診斷",
            "年度行銷目標 KPI 設定",
            "媒體策略全漏斗規劃",
            "整合媒體策略建議",
            "預算配置方案",
            "Meta / LINE 廣告策略",
            "Google / SEO 搜尋策略",
            "YouTube 影音行銷策略",
            "KOL / 網紅行銷規劃",
            "LINE CRM 會員經營策略",
            "季度執行時程計畫",
            "下一步行動建議",
        ]
    topics_str = "\n".join(f"{i+1}. {t}" for i, t in enumerate(topics))
    n_topics = len(topics)

    prompt = f"""你是潮網科技的業務顧問，正在為客戶製作提案簡報。

【重要原則】
- 所有內容必須根據下方提供的資料撰寫，不得憑空捏造
- 若官網摘要有提供，以官網內容為主要依據
- 不要假設公司的產品或服務，必須從資料中找依據
{extra_section}{website_section}
【廠商基本資料】
- 公司：{lead.company_name}
- 產業：{lead.industry or "（請從官網內容判斷）"}
- 城市：{lead.city or "台灣"}
- 規模：{lead.company_size or "未知"}
- 含金量：{score}/100
- 數位追蹤工具：{tech_str}
- 廣告投放：{ad_str}

【簡報頁面主題（依序生成 {n_topics} 頁，每頁對應一個主題）】
{topics_str}

只輸出 JSON，不要其他說明：
{{"slides": [{{"title": "標題", "bullets": ["要點1", "要點2", "要點3", "要點4", "要點5"]}}]}}

嚴格要求：
- 必須輸出 {n_topics} 個 slides，每頁對應上方一個主題
- 每頁 4-5 個繁體中文條列，內容必須具體且與該公司實際業務相符
- 數字、百分比、策略工具等要具體，不要用「例如」或模板文字"""

    try:
        from openai import OpenAI
        oai = OpenAI(api_key=OPENAI_API_KEY)
        if body.context_images:
            content: list = [{"type": "text", "text": prompt}]
            for img_url in body.context_images[:4]:  # max 4 images
                content.append({"type": "image_url", "image_url": {"url": img_url, "detail": "low"}})
            messages = [{"role": "user", "content": content}]
        else:
            messages = [{"role": "user", "content": prompt}]
        resp = oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            response_format={"type": "json_object"},
        )
        slide_content = json.loads(resp.choices[0].message.content.strip())
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 內容生成失敗：{str(e)}")

    return {
        "company_name": lead.company_name,
        "industry": lead.industry or "",
        "slides": slide_content.get("slides", []),
    }
