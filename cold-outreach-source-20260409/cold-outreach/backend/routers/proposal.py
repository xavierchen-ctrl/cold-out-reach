import asyncio
import io
import json
import os
import re
import time
import uuid
import httpx
from concurrent.futures import ThreadPoolExecutor
from urllib.parse import quote
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional
from openai import OpenAI
from sqlalchemy.orm import Session
from auth import get_current_user
from database import get_db
from models import User, Lead

router = APIRouter(prefix="/api/proposal", tags=["proposal"])
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")

# ── Assistants API job store ───────────────────────────────────────────────────
_jobs: dict = {}          # job_id → {status, message, bytes, filename, created_at}
_executor = ThreadPoolExecutor(max_workers=3)
_assistant_id: str | None = None

_HTTP_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
        "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120 Safari/537.36"
    )
}


def _strip_tags(html: str) -> str:
    text = re.sub(r'<[^>]+>', ' ', html)
    return re.sub(r'\s+', ' ', text).strip()

# ── Colors ────────────────────────────────────────────────────────────────────
C_NAVY   = (0x1B, 0x3A, 0x6B)
C_BLUE   = (0x21, 0x96, 0xF3)
C_LBLUE  = (0xE3, 0xF2, 0xFD)
C_ORANGE = (0xF5, 0x7C, 0x00)
C_LORANGE= (0xFF, 0xF3, 0xE0)
C_WHITE  = (0xFF, 0xFF, 0xFF)
C_LGRAY  = (0xF5, 0xF5, 0xF5)
C_DGRAY  = (0x42, 0x42, 0x42)
C_MGRAY  = (0x75, 0x75, 0x75)
C_BORDER = (0xDD, 0xDD, 0xDD)
C_LBLUE2 = (0xBB, 0xDE, 0xFB)
C_GREEN  = (0x2E, 0x7D, 0x32)
C_LGREEN = (0xE8, 0xF5, 0xE9)


class ProposalRequest(BaseModel):
    client_name: str
    industry: str
    current_situation: str
    services: List[str]
    monthly_budget: str
    special_notes: Optional[str] = None
    year: int = 2026


@router.post("/generate")
async def generate_proposal(
    body: ProposalRequest,
    current_user: User = Depends(get_current_user),
):
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    try:
        content = await _generate_content_unified(body)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 生成失敗：{str(e)}")

    try:
        pptx_bytes = _build_pptx(content, body)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX 建立失敗：{str(e)}")

    filename = f"{body.client_name}_{body.year}_媒體提案.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(filename)}"},
    )


class GenerateFromLeadRequest(BaseModel):
    lead_id: str
    services: List[str] = ["廣告投放", "SEO優化"]
    budget_range: str = "50-100萬/月"
    extra_context: str = ""
    year: int = 2026


_BUDGET_MAP = {
    "10-30萬/月": "30",
    "30-50萬/月": "50",
    "50-100萬/月": "100",
    "100萬以上/月": "150",
}


def _get_or_create_assistant(client: OpenAI) -> str:
    global _assistant_id
    if _assistant_id:
        try:
            client.beta.assistants.retrieve(_assistant_id)
            return _assistant_id
        except Exception:
            pass
    assistant = client.beta.assistants.create(
        name="潮網科技 提案設計師",
        instructions=(
            "你是潮網科技的資深數位行銷顧問兼專業簡報設計師。"
            "當收到提案要求時，請用 python-pptx 製作視覺豐富、每頁版面皆不同的專業簡報。"
            "程式碼必須完整可執行，存檔路徑為 /tmp/proposal.pptx。"
        ),
        tools=[{"type": "code_interpreter"}],
        model="gpt-4o",
    )
    _assistant_id = assistant.id
    return _assistant_id


def _build_assistants_prompt(client_name: str, industry: str, budget: str,
                              services: str, current_situation: str, year: int) -> str:
    return f"""你是潮網科技的資深數位行銷顧問兼專業簡報設計師。

【客戶資訊】
客戶名稱：{client_name}
產業：{industry}
月預算：{budget}萬
主推服務：{services}
品牌現況：{current_situation}

請完成以下任務：
1. 根據客戶資訊生成完整的媒體提案內容（繁體中文）
2. 用 python-pptx 製作一份專業的 {year} 年度數位行銷媒體提案 PowerPoint（16頁）
3. 存檔為 /tmp/proposal.pptx

【技術規格】
- 簡報尺寸：13.33 英吋 × 7.5 英吋（16:9）
- python-pptx 1.0.x：所有 EMU 值必須用 int() 轉換，shp.line.width = 0 表示無邊框
- 使用 Presentation()、slide_layouts[6]（空白版型）

【設計要求】
- 每頁使用不同的版面佈局（嚴禁每頁設計相同）
- 主色系：深藍 RGB(27,58,107) + 橘色 RGB(245,124,0) + 白色
- 每頁都要有色塊、標題、內容；善用對角切版、左右分欄、卡片格狀等設計
- 重要數據用大字號突出顯示
- 所有文字使用繁體中文

【頁面結構】
第1頁：封面 — 深藍背景、客戶名稱大標、「{year} 年度數位行銷媒體提案」副標、潮網科技
第2頁：品牌現況分析 — 3個優勢 + 現況說明
第3頁：目標客群分析 — 3個族群卡片（年齡、需求、決策關鍵）
第4頁：行銷問題診斷 — 3個問題診斷卡片
第5頁：年度 KPI 目標 — 5個指標，大數字視覺化
第6頁：全漏斗媒體策略 — Awareness → Consideration → Conversion → Retention 流程
第7頁：整合媒體策略總表 — 表格（漏斗階段 / 媒體工具 / 溝通核心）
第8頁：預算配置 — 條狀圖（Meta 30%、Google 20%、YouTube 10%、LinkedIn 15%、LINE 10%、KOL 10%、展會 5%）
第9頁：Meta/Facebook/Instagram 廣告策略 — 3個受眾族群
第10頁：Google/SEO 策略 — 3種關鍵字類型
第11頁：YouTube 影音策略 — 3種內容角色
第12頁：KOL/網紅行銷規劃 — Tier 1/2/3 三層架構
第13頁：LinkedIn B2B 策略 — 目標職稱、廣告形式
第14頁：LINE CRM 會員旅程 — Day 0 到 Day 30+ 時間軸
第15頁：季度執行計畫 — Q1~Q4 橫向時程表
第16頁：結語 — 深藍背景、總結訊息、下一步行動

請直接撰寫並執行完整的 Python 程式碼（無需解釋），確認 /tmp/proposal.pptx 已成功儲存。"""


def _do_assistants_job(job_id: str, prompt: str, filename: str) -> None:
    try:
        _jobs[job_id].update(status="running", message="AI 正在初始化...")
        client = OpenAI(api_key=OPENAI_API_KEY)
        assistant_id = _get_or_create_assistant(client)

        thread = client.beta.threads.create()
        client.beta.threads.messages.create(
            thread_id=thread.id, role="user", content=prompt
        )

        run = client.beta.threads.runs.create(
            thread_id=thread.id, assistant_id=assistant_id
        )

        # Poll with status updates
        while run.status in ("queued", "in_progress", "cancelling"):
            time.sleep(6)
            run = client.beta.threads.runs.retrieve(
                thread_id=thread.id, run_id=run.id
            )
            status_map = {
                "queued": "AI 正在排隊中...",
                "in_progress": "AI 正在撰寫程式碼並設計簡報，請稍候...",
                "cancelling": "正在取消...",
            }
            _jobs[job_id]["message"] = status_map.get(run.status, run.status)

        if run.status != "completed":
            last_error = getattr(run, "last_error", None)
            err_msg = str(last_error) if last_error else run.status
            _jobs[job_id].update(status="error", message=f"生成失敗：{err_msg}")
            return

        # Find the PPTX file attachment
        messages = client.beta.threads.messages.list(thread_id=thread.id)
        pptx_bytes = None
        for msg in messages.data:
            if msg.role != "assistant":
                continue
            for attachment in (getattr(msg, "attachments", None) or []):
                fid = getattr(attachment, "file_id", None)
                if fid:
                    pptx_bytes = client.files.content(fid).read()
                    break
            if pptx_bytes:
                break

        if not pptx_bytes:
            _jobs[job_id].update(status="error", message="找不到生成的 PPTX 檔案，請重試")
            return

        _jobs[job_id].update(
            status="done",
            message="簡報已完成！",
            pptx_bytes=pptx_bytes,
            filename=filename,
            done_at=time.time(),
        )

    except Exception as e:
        _jobs[job_id].update(status="error", message=str(e))


@router.post("/chatgpt-prompt")
async def get_chatgpt_prompt(
    body: GenerateFromLeadRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    monthly_budget = _BUDGET_MAP.get(body.budget_range, "100")

    website_summary = ""
    if lead.website:
        try:
            url = lead.website if lead.website.startswith("http") else f"https://{lead.website}"
            async with httpx.AsyncClient(timeout=6, follow_redirects=True, headers=_HTTP_HEADERS) as hc:
                res = await hc.get(url)
                website_summary = _strip_tags(res.text)[:1000]
        except Exception:
            pass

    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    tech_str = "、".join(k for k, v in [("GA4", tech.get("ga4")), ("GTM", tech.get("gtm")), ("Meta Pixel", tech.get("meta_pixel"))] if v)
    ad_str = "、".join(k for k, v in [("Meta 廣告", ad.get("meta", {}).get("has_ads")), ("Google 廣告", ad.get("google_ads", {}).get("has_ads"))] if v)

    parts = []
    if lead.industry:      parts.append(f"產業：{lead.industry}")
    if lead.company_size:  parts.append(f"規模：{lead.company_size}")
    if lead.city:          parts.append(f"城市：{lead.city}")
    if tech_str:           parts.append(f"數位工具：{tech_str}")
    if ad_str:             parts.append(f"廣告現況：{ad_str}")
    if website_summary:    parts.append(f"官網：{website_summary}")
    if body.extra_context.strip(): parts.append(f"補充：{body.extra_context.strip()}")
    current_situation = "\n".join(parts) or f"{lead.company_name} 的品牌與行銷現況"

    services_str = "、".join(body.services)

    budget_int = int(monthly_budget)
    if budget_int <= 30:
        budget_guidance = """【預算策略指引 — 小預算（30萬以下）】
月預算有限，建議聚焦 2-3 個核心管道，不宜分散：
- 必選：Meta FB/IG 廣告（品牌+效果兼顧）、Google Search（捕捉主動需求）
- 可視情況加入：LINE OA 基礎經營
- 暫緩：YouTube 大規模投放、KOL、LinkedIn、展會（等預算擴大再納入）
- 重點：把有限預算集中在轉換率最高的管道，追求 ROAS 最大化"""
    elif budget_int <= 50:
        budget_guidance = """【預算策略指引 — 中小預算（30-50萬）】
建議 3-4 個管道組合，兼顧品效：
- 核心：Meta 廣告 + Google Ads
- 擴充：可加入 KOL 微網紅（口碑擴散）或 LINE CRM（會員經營）擇一
- 可視需求加：YouTube 短影音（不做大規模投放）
- 暫緩：LinkedIn、大型展會"""
    elif budget_int <= 100:
        budget_guidance = """【預算策略指引 — 中大預算（50-100萬）】
建議 4-5 個管道全數位整合：
- 核心：Meta + Google + YouTube 影音廣告
- 加入：KOL 行銷（口碑+觸及）、LINE CRM 自動化（會員留存）
- 視產業：若有 B2B 需求可加 LinkedIn 廣告
- 暫緩：大型展會、Podcast 冠名（除非產業特別適合）"""
    else:
        budget_guidance = """【預算策略指引 — 大預算（100萬以上）】
全管道整合，建立品牌資產與多觸點佈局：
- 全數位：Meta + Google + YouTube + KOL + LinkedIn + LINE CRM
- 線下延伸：展會、研討會、Podcast 冠名、媒體原生廣告
- 重點：各管道設定明確 KPI（品牌曝光 vs 效果轉換），避免預算重疊浪費
- 建議設立整合儀表板，跨管道歸因分析"""

    prompt = f"""請幫我製作一份專業的數位行銷媒體提案 PowerPoint（產生可下載的 .pptx 檔案，繁體中文）。

【客戶資訊】
客戶名稱：{lead.company_name}
產業：{lead.industry or "未知"}
月預算：{monthly_budget}萬
主推服務：{services_str}
品牌現況：
{current_situation}

{budget_guidance}

【簡報結構（頁碼固定，請嚴格遵守）】

第 1 頁：客製化封面
- 主標：{lead.company_name} × 潮網科技
- 副標：數位行銷媒體提案
- 年度標示（2025 / 2026）

--- 以下第 2-10 頁為固定公司介紹，內容固定，請完整呈現 ---

第 2 頁：OUR HISTORY — 潮網的成長歷程
- 2010：潮網的起點（數位廣告代理起家）
- 2012：進軍國際（佈局海外市場）
- 2014：走入 ADTECH（自建廣告技術平台 PUNWAVE）
- 2018：佈局 MARTECH（推出 ORCA CDP & MA+）
- 2025：持續建構行銷生態圈（整合 AI、數據、媒體全鏈路）

第 3 頁：SERVICE — 數據驅動行銷，四大服務類別
「數位廣告」
  核心平台：LUCA 一站式整合平台、Narwhal
  媒體渠道：Meta / Google / YouTube / LinkedIn / LINE / Criteo / X / Yahoo / Dcard / The Trade Desk
「整合行銷」
  數位行銷、口碑行銷、內容合作、網紅行銷、跨境行銷
  社群行銷、輿情分析、公關策略、線上&線下活動
「電子商務」
  平台：Shopify / SHOPLINE / GA4
  服務：平台導入教育訓練、銷售網站建置優化、代碼串接埋設建置、GA4移轉帳戶健檢、數據分析指標管理
「行銷科技」
  核心產品：ORCA CDP & ORCA MA+
  服務：ORCA CDP導入應用、ORCA MA+再行銷運用、BONITO輿情分析平台、Chatbot串接運營策略、跨裝置行為數據媒合
  合作夥伴：Salesforce / Oracle

第 4 頁：PUNWAVE AI — 跨媒體採購 × 人工智慧
- 以 AI 人工智慧、大數據分析為核心
- 系統節點：PUNWAVE RTB（即時競價）→ PUNWAVE DMP（資料整合/標籤分類/分析）→ PUNWAVE ATD → PUNWAVE AI Smart Manager（24小時即時優化動態預算分配）→ 串接 Facebook / Google / PUNWAVE DSP

第 5 頁：Excel 報表範例（每日數據）
- 每日廣告數據追蹤報表
- 欄位：日期 / 廣告策略 / 曝光數 / 點擊數 / CTR / CPC / 成本 / 花費
- 涵蓋多種廣告策略（品牌曝光、再行銷、關鍵字）的每日數據

第 6 頁：Excel 報表範例（媒體策略分析）
- 媒體數據策略表現彙整（Meta / Google / YouTube / LINE）
- 顧問 Finding 分析：當期廣告洞察、優化建議、下期策略調整方向

第 7 頁：圖像報表範例（一）
- 客製化圖像分析報表，包含：
  ① TA 分析（受眾年齡/性別/興趣分佈）
  ② 逐週表現（曝光、點擊、轉換週別趨勢）
  ③ 競價分析（CPM/CPC 與產業均值對比）

第 8 頁：圖像報表範例（二）
- 多媒體整合成效報表，包含：
  ① 各廣告策略成效比較
  ② Google 關鍵字月別表現
  ③ Criteo 廣告成效（再行銷 ROAS 與轉換數趨勢）

第 9 頁：雲端素材表範例
- Meta 廣告素材雲端管理表（欄位：貼文連結 / 廣告目標 / 格式 / 上線日期 / 狀態）
- 8 種 Meta 廣告格式：單圖 / 輪播 / 影片 / 限時動態 / Reels / 即時體驗 / 精選集 / 動態廣告

第 10 頁：為什麼選擇潮網？— 4 大優勢 × 8 大特色
4 大優勢：
  ① 成效導向：以 KPI 與 ROAS 為核心，追求實際業績成長
  ② 精準投放：結合 AI 智能與大數據，精準觸達目標受眾
  ③ 數據分析：深度數據洞察，每日追蹤與即時優化
  ④ 全渠道整合服務：線上+線下、品牌+效果，一站式整合
8 大特色：
  ✦ AI 智能分析精準鎖定目標受眾
  ✦ 精準投放並成效追蹤
  ✦ 專屬客製化行銷數據分析＆報表策略
  ✦ 數位廣告 × 內容行銷線上與線下整合
  ✦ Daily 確認廣告狀況，即時調整
  ✦ 產業經驗豐富的分析師、企劃、編輯等策略制定
  ✦ 協助 GA 數據追蹤，精準數位分析交叉比對
  ✦ 會員經營 × 數據整合，打造品牌增長飛輪

--- 以下第 11 頁起為客製化提案內容 ---

【客製化提案內容方向（第 11 頁起，依客戶情況取捨）】
- 品牌現況分析
- 目標客群分析
- 行銷問題診斷
- 年度 KPI 目標
- 全漏斗媒體策略
- 整合媒體策略總表
- 預算配置
- 各媒體管道策略（依上方預算策略指引，只納入適合的管道，每個管道獨立一頁說明策略與執行方式）
- 季度執行計畫
- 結語與下一步行動

請根據客戶的產業、預算規模與主推服務，自行判斷哪些內容值得獨立一頁、哪些可以合併，產生最合適的頁數。
請產生完整可下載的 PPTX 檔案。"""

    return {"prompt": prompt}


@router.post("/generate-ai-from-lead")
async def generate_ai_from_lead(
    body: GenerateFromLeadRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    monthly_budget = _BUDGET_MAP.get(body.budget_range, "100")

    # Scrape website for context
    website_summary = ""
    if lead.website:
        try:
            url = lead.website if lead.website.startswith("http") else f"https://{lead.website}"
            async with httpx.AsyncClient(timeout=8, follow_redirects=True, headers=_HTTP_HEADERS) as hc:
                res = await hc.get(url)
                website_summary = _strip_tags(res.text)[:1500]
        except Exception:
            pass

    # Build current_situation
    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    tech_str = "、".join(k for k, v in [("GA4", tech.get("ga4")), ("GTM", tech.get("gtm")), ("Meta Pixel", tech.get("meta_pixel"))] if v)
    ad_str = "、".join(k for k, v in [("Meta 廣告", ad.get("meta", {}).get("has_ads")), ("Google 廣告", ad.get("google_ads", {}).get("has_ads"))] if v)

    parts = []
    if lead.industry:   parts.append(f"產業：{lead.industry}")
    if lead.company_size: parts.append(f"規模：{lead.company_size}")
    if lead.city:       parts.append(f"城市：{lead.city}")
    if tech_str:        parts.append(f"數位工具：{tech_str}")
    if ad_str:          parts.append(f"廣告現況：{ad_str}")
    if website_summary: parts.append(f"官網：{website_summary}")
    if body.extra_context.strip(): parts.append(f"補充：{body.extra_context.strip()}")
    current_situation = "\n".join(parts) or f"{lead.company_name} 的品牌與行銷現況"

    services_str = "、".join(body.services)
    prompt = _build_assistants_prompt(
        client_name=lead.company_name,
        industry=lead.industry or "未知",
        budget=monthly_budget,
        services=services_str,
        current_situation=current_situation,
        year=body.year,
    )

    job_id = str(uuid.uuid4())
    filename = f"{lead.company_name}_{body.year}_媒體提案.pptx"
    _jobs[job_id] = {
        "status": "queued",
        "message": "已加入佇列，即將開始...",
        "filename": filename,
        "created_at": time.time(),
    }

    loop = asyncio.get_event_loop()
    loop.run_in_executor(_executor, _do_assistants_job, job_id, prompt, filename)

    return {"job_id": job_id}


@router.get("/job/{job_id}")
async def get_job_status(
    job_id: str,
    current_user: User = Depends(get_current_user),
):
    job = _jobs.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return {"status": job["status"], "message": job["message"]}


@router.get("/job/{job_id}/file")
async def download_job_file(
    job_id: str,
    current_user: User = Depends(get_current_user),
):
    job = _jobs.get(job_id)
    if not job or job["status"] != "done":
        raise HTTPException(status_code=404, detail="File not ready")
    pptx_bytes = job["pptx_bytes"]
    filename = job["filename"]
    _jobs.pop(job_id, None)
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(filename)}"},
    )


@router.post("/generate-from-lead")
async def generate_from_lead(
    body: GenerateFromLeadRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OpenAI API key not configured")

    lead = db.query(Lead).filter(Lead.id == body.lead_id).first()
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    monthly_budget = _BUDGET_MAP.get(body.budget_range, "100")

    # Scrape website for context
    website_summary = ""
    if lead.website:
        try:
            url = lead.website if lead.website.startswith("http") else f"https://{lead.website}"
            async with httpx.AsyncClient(timeout=10, follow_redirects=True, headers=_HTTP_HEADERS) as client:
                res = await client.get(url)
                website_summary = _strip_tags(res.text)[:2000]
        except Exception:
            pass

    # Build current_situation from enriched lead data
    tech = lead.tech_signals or {}
    ad = lead.ad_signals or {}
    tech_str = "、".join(k for k, v in [("GA4", tech.get("ga4")), ("GTM", tech.get("gtm")), ("Meta Pixel", tech.get("meta_pixel"))] if v)
    ad_str = "、".join(k for k, v in [("Meta 廣告", ad.get("meta", {}).get("has_ads")), ("Google 廣告", ad.get("google_ads", {}).get("has_ads"))] if v)
    score = lead.enriched_score or 0

    parts = []
    if lead.industry:
        parts.append(f"產業類別：{lead.industry}")
    if lead.company_size:
        parts.append(f"公司規模：{lead.company_size}")
    if lead.city:
        parts.append(f"所在城市：{lead.city}")
    if tech_str:
        parts.append(f"已使用數位追蹤工具：{tech_str}")
    if ad_str:
        parts.append(f"目前廣告投放：{ad_str}")
    if score:
        parts.append(f"含金量評分：{score}/100")
    if website_summary:
        parts.append(f"官網內容：{website_summary}")
    if body.extra_context.strip():
        parts.append(f"補充說明：{body.extra_context.strip()}")

    current_situation = "\n".join(parts) or f"{lead.company_name} 的品牌與行銷現況"

    proposal_req = ProposalRequest(
        client_name=lead.company_name,
        industry=lead.industry or "未知",
        current_situation=current_situation,
        services=body.services,
        monthly_budget=monthly_budget,
        special_notes=None,
        year=body.year,
    )

    try:
        content = await _generate_content_unified(proposal_req)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI 生成失敗：{str(e)}")

    try:
        pptx_bytes = _build_pptx(content, proposal_req)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX 建立失敗：{str(e)}")

    filename = f"{lead.company_name}_{body.year}_媒體提案.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(filename)}"},
    )


# ── AI content generation ─────────────────────────────────────────────────────

async def _generate_content_unified(body: ProposalRequest) -> dict:
    services_str = "、".join(body.services)
    special = f"\n特殊需求：{body.special_notes}" if body.special_notes else ""
    budget_num = re.sub(r"[^\d]", "", body.monthly_budget) or "100"

    prompt = f"""你是潮網科技的資深數位行銷顧問，請根據以下資訊生成完整媒體提案內容（繁體中文）。

客戶：{body.client_name} | 產業：{body.industry}
月預算：{budget_num}萬 | 主推服務：{services_str}
品牌現況：{body.current_situation}{special}

本提案涵蓋所有行銷管道（Meta、Google、YouTube、KOL、LinkedIn、LINE CRM、展會、內容資產等），
客戶可依實際需求自行刪減不適用的部分。請根據客戶產業特性生成最貼切的內容（勿使用預設範例）。

輸出嚴格 JSON：

{{
  "subtitle": "副標題（20字內，描述品牌/業務成長方向）",
  "brand_strengths": ["品牌/產品優勢1", "品牌/產品優勢2", "品牌/產品優勢3"],
  "brand_d2c": "品牌直客/通路/商業模式現況說明（40字內）",
  "market_segments": [
    {{"age": "族群特徵1", "name": "族群名稱", "needs": ["需求1", "需求2"], "decision": "決策關鍵"}},
    {{"age": "族群特徵2", "name": "族群名稱", "needs": ["需求1", "需求2"], "decision": "決策關鍵"}},
    {{"age": "族群特徵3", "name": "族群名稱", "needs": ["需求1", "需求2"], "decision": "決策關鍵"}}
  ],
  "problems": [
    {{"title": "問題標題1", "desc": "問題說明（30字內）"}},
    {{"title": "問題標題2", "desc": "問題說明"}},
    {{"title": "問題標題3", "desc": "問題說明"}}
  ],
  "kpis": [
    {{"label": "KPI名稱1", "value": "+50%"}},
    {{"label": "KPI名稱2", "value": "+80%"}},
    {{"label": "KPI名稱3", "value": "+40%"}},
    {{"label": "KPI名稱4", "value": "4.0+"}},
    {{"label": "KPI名稱5", "value": "+20%"}}
  ],
  "media_strategy": [
    {{"stage": "認知 (Awareness)", "tools": "媒體工具（逗號分隔）", "message": "溝通核心（20字）"}},
    {{"stage": "考慮 (Consideration)", "tools": "媒體工具", "message": "溝通核心"}},
    {{"stage": "轉換 (Conversion)", "tools": "媒體工具", "message": "溝通核心"}},
    {{"stage": "回購/留存 (Retention)", "tools": "媒體工具", "message": "溝通核心"}}
  ],
  "meta_audiences": [
    {{"name": "受眾族群1", "products": "主打商品/服務", "creative": "廣告創意方向（15字）"}},
    {{"name": "受眾族群2", "products": "主打商品/服務", "creative": "廣告創意方向"}},
    {{"name": "受眾族群3", "products": "主打商品/服務", "creative": "廣告創意方向"}}
  ],
  "google_keywords": [
    {{"type": "關鍵字類型1", "keywords": "關鍵字1、關鍵字2、關鍵字3", "effect": "效益說明"}},
    {{"type": "關鍵字類型2", "keywords": "關鍵字1、關鍵字2、關鍵字3", "effect": "效益說明"}},
    {{"type": "關鍵字類型3", "keywords": "品牌詞1、品牌詞2", "effect": "效益說明"}}
  ],
  "youtube_experts": [
    {{"role": "內容角色1", "content": "內容方向（20字）"}},
    {{"role": "內容角色2", "content": "內容方向"}},
    {{"role": "內容角色3", "content": "內容方向"}}
  ],
  "kol_tiers": [
    {{"tier": "Tier 1  大型KOL", "purpose": "流量爆發", "desc": "25字以內說明"}},
    {{"tier": "Tier 2  專業人士", "purpose": "信任轉化", "desc": "25字以內說明"}},
    {{"tier": "Tier 3  微網紅(KOC)", "purpose": "社群擴散", "desc": "25字以內說明"}}
  ],
  "linkedin_targeting": {{
    "roles": ["目標職稱1", "目標職稱2", "目標職稱3"],
    "companies": ["目標公司類型1", "目標公司類型2"],
    "ad_formats": ["廣告形式1", "廣告形式2"]
  }},
  "seo_keywords": [
    {{"type": "SEO類型1", "keywords": "關鍵字1、關鍵字2", "intent": "搜尋意圖"}},
    {{"type": "SEO類型2", "keywords": "關鍵字1、關鍵字2", "intent": "搜尋意圖"}},
    {{"type": "SEO類型3", "keywords": "品牌詞1、品牌詞2", "intent": "品牌認知"}}
  ],
  "thought_leadership": [
    {{"format": "內容形式1", "topic": "主題（20字）", "goal": "目標（15字）"}},
    {{"format": "內容形式2", "topic": "主題", "goal": "目標"}},
    {{"format": "內容形式3", "topic": "主題", "goal": "目標"}}
  ],
  "events": [
    {{"name": "活動/展會名稱1", "strategy": "參與策略", "tactic": "具體戰術"}},
    {{"name": "活動/展會名稱2", "strategy": "參與策略", "tactic": "具體戰術"}}
  ],
  "crm_steps": [
    {{"day": "Day 0",   "title": "初始接觸",   "desc": "行動說明"}},
    {{"day": "Day 1-7", "title": "培育互動",   "desc": "行動說明"}},
    {{"day": "Day 14",  "title": "深化關係",   "desc": "行動說明"}},
    {{"day": "Day 25",  "title": "促進轉換",   "desc": "行動說明"}},
    {{"day": "Day 30+", "title": "持續回購/合作", "desc": "行動說明"}}
  ],
  "content_assets": [
    {{"type": "內容資產類型1", "desc": "說明（25字）"}},
    {{"type": "內容資產類型2", "desc": "說明"}},
    {{"type": "內容資產類型3", "desc": "說明"}}
  ],
  "must_buy": ["必購媒體1", "必購媒體2", "必購媒體3", "必購媒體4"],
  "bonus_resources": ["加值資源1", "加值資源2", "加值資源3"],
  "quarterly_plan": [
    {{"quarter": "Q1  策略方向", "goal": "目標說明", "strategy": "執行策略"}},
    {{"quarter": "Q2  策略方向", "goal": "目標說明", "strategy": "執行策略"}},
    {{"quarter": "Q3  策略方向", "goal": "目標說明", "strategy": "執行策略"}},
    {{"quarter": "Q4  策略方向", "goal": "目標說明", "strategy": "執行策略"}}
  ],
  "closing_message": "結語（80字以內，說明整合行銷策略的核心價值與預期成效）"
}}"""

    client = OpenAI(api_key=OPENAI_API_KEY)
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    )
    return json.loads(response.choices[0].message.content)


# ── PPTX helpers ──────────────────────────────────────────────────────────────

def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(*rgb)


def _rect(slide, l, t, w, h, fill=None, border=None, border_w=12700):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(*fill)
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = _rgb(*border)
        shp.line.width = Emu(border_w)
    else:
        shp.line.width = 0
    return shp


def _txt(slide, text, l, t, w, h, size=14, bold=False, color=None, align="left", italic=False, wrap=True):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    color = color or C_DGRAY
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    box = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(*color)
    return box, tf


def _bullets(slide, items, l, t, w, h, size=12, color=None, title=None, title_size=15, title_color=None, title_bold=True):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    color = color or C_DGRAY
    title_color = title_color or C_NAVY
    box = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = title_bold
        run.font.color.rgb = _rgb(*title_color)
        first = False
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"· {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(*color)
    return box


def _header(slide, title):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    _rect(slide, int(Inches(0)), int(Inches(0)), int(Inches(13.33)), int(Inches(1.05)), fill=C_NAVY)
    _rect(slide, int(Inches(0.35)), int(Inches(0.83)), int(Inches(1.5)), int(Inches(0.05)), fill=C_BLUE)
    box = slide.shapes.add_textbox(int(Inches(0.35)), int(Inches(0.1)), int(Inches(12.5)), int(Inches(0.72)))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _rgb(*C_WHITE)


def _table(slide, headers, rows, l, t, w, h):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    ncols = len(headers)
    nrows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(nrows, ncols, int(l), int(t), int(w), int(h))
    tbl = tbl_shape.table
    col_w = int(w) // ncols
    for i in range(ncols):
        tbl.columns[i].width = col_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(*C_NAVY)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = _rgb(*C_WHITE)

    # Data rows
    for ri, row in enumerate(rows):
        bg = C_LGRAY if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(*bg)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci > 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(12)
            run.font.bold = (ci == 0)
            run.font.color.rgb = _rgb(*C_NAVY if ci == 0 else C_DGRAY)

    return tbl


# ── B2C Individual slides ─────────────────────────────────────────────────────

def _slide_cover(prs, layout, client_name, subtitle, year):
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        slide = prs.slides.add_slide(layout)
        _set_bg(slide, C_NAVY)
        _rect(slide, int(Inches(0)), int(Inches(3.0)), int(Inches(13.33)), int(Inches(0.08)), fill=C_BLUE)
        _txt(slide, f"{client_name} {year}", int(Inches(1.5)), int(Inches(1.2)), int(Inches(10)), int(Inches(1.2)),
             size=40, bold=True, color=C_WHITE, align="center")
        _txt(slide, "數位媒體成長提案", int(Inches(1.5)), int(Inches(2.2)), int(Inches(10)), int(Inches(0.9)),
             size=30, bold=True, color=C_BLUE, align="center")
        _txt(slide, subtitle, int(Inches(2)), int(Inches(3.3)), int(Inches(9.33)), int(Inches(0.7)),
             size=16, color=C_WHITE, align="center", italic=True)
        _txt(slide, "Powered by 潮網科技 Wavenet Technology", int(Inches(3)), int(Inches(6.8)), int(Inches(7.33)), int(Inches(0.4)),
             size=11, color=C_MGRAY, align="center")
    except Exception as e:
        print(f"[slide_cover error] {e}")


def _slide_brand(prs, layout, strengths, d2c_desc):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "品牌現況分析")
        _rect(slide, int(Inches(0.4)), int(Inches(1.25)), int(Inches(5.8)), int(Inches(5.5)), fill=C_LGRAY, border=C_BORDER)
        _bullets(slide, strengths, int(Inches(0.6)), int(Inches(1.45)), int(Inches(5.4)), int(Inches(5.0)),
                 title="商品線完整", title_size=16, size=14)
        _rect(slide, int(Inches(7.1)), int(Inches(1.25)), int(Inches(5.8)), int(Inches(5.5)), fill=C_LBLUE, border=C_BORDER)
        _txt(slide, "具備D2C經營基礎", int(Inches(7.3)), int(Inches(1.5)), int(Inches(5.4)), int(Inches(0.5)),
             size=16, bold=True, color=C_NAVY)
        _txt(slide, d2c_desc, int(Inches(7.3)), int(Inches(2.1)), int(Inches(5.4)), int(Inches(2.0)),
             size=13, color=C_DGRAY, wrap=True)
        _rect(slide, int(Inches(7.3)), int(Inches(5.0)), int(Inches(5.4)), int(Inches(0.6)), fill=C_LBLUE2)
        _txt(slide, "核心優勢：數據自主性高、會員導購潛力大。", int(Inches(7.5)), int(Inches(5.05)), int(Inches(5.0)), int(Inches(0.5)),
             size=12, bold=True, color=C_NAVY)
    except Exception as e:
        print(f"[slide_brand error] {e}")


def _slide_market(prs, layout, segments, year):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, f"{year} 市場三大成長族群")
        card_w = int(Inches(3.8))
        card_h = int(Inches(5.3))
        xs = [int(Inches(0.4)), int(Inches(4.76)), int(Inches(9.12))]
        for i, seg in enumerate(segments[:3]):
            x = xs[i]
            _rect(slide, x, int(Inches(1.25)), card_w, card_h, fill=C_WHITE, border=C_BORDER)
            _txt(slide, seg.get("age", ""), x, int(Inches(1.4)), card_w, int(Inches(0.4)),
                 size=12, color=C_MGRAY, align="center")
            _txt(slide, seg.get("name", ""), x, int(Inches(1.8)), card_w, int(Inches(0.5)),
                 size=16, bold=True, color=C_NAVY, align="center")
            needs = seg.get("needs", [])
            _bullets(slide, needs, x + int(Inches(0.2)), int(Inches(2.4)), card_w - int(Inches(0.4)), int(Inches(2.0)),
                     title="需求：", title_size=13, size=13)
            _txt(slide, f"決策關鍵：{seg.get('decision', '')}", x + int(Inches(0.2)), int(Inches(4.9)),
                 card_w - int(Inches(0.4)), int(Inches(0.5)), size=12, bold=True, color=C_ORANGE)
    except Exception as e:
        print(f"[slide_market error] {e}")


def _slide_problems(prs, layout, problems, title="品牌問題診斷"):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, title)
        _txt(slide, "⚠", int(Inches(0.5)), int(Inches(2.5)), int(Inches(1.5)), int(Inches(2.0)),
             size=60, color=C_ORANGE, align="center")
        _txt(slide, "當前增長瓶頸", int(Inches(0.3)), int(Inches(4.7)), int(Inches(1.8)), int(Inches(0.4)),
             size=12, color=C_MGRAY, align="center")
        for i, prob in enumerate(problems[:4]):
            y = int(Inches(1.35)) + i * int(Inches(1.35))
            _rect(slide, int(Inches(2.3)), y, int(Inches(10.6)), int(Inches(1.15)), fill=C_LORANGE, border=C_BORDER)
            _txt(slide, prob.get("title", ""), int(Inches(2.55)), y + int(Inches(0.1)), int(Inches(10.0)), int(Inches(0.4)),
                 size=15, bold=True, color=C_ORANGE)
            _txt(slide, prob.get("desc", ""), int(Inches(2.55)), y + int(Inches(0.5)), int(Inches(10.0)), int(Inches(0.55)),
                 size=13, color=C_DGRAY)
    except Exception as e:
        print(f"[slide_problems error] {e}")


def _slide_kpis(prs, layout, kpis, year):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, f"{year} 數位媒體年度 KPI")
        positions = [
            (int(Inches(0.5)),  int(Inches(1.8))),
            (int(Inches(4.5)),  int(Inches(1.8))),
            (int(Inches(8.5)),  int(Inches(1.8))),
            (int(Inches(2.3)),  int(Inches(4.2))),
            (int(Inches(6.3)),  int(Inches(4.2))),
        ]
        for i, kpi in enumerate(kpis[:5]):
            x, y = positions[i]
            _txt(slide, kpi.get("value", ""), x, y, int(Inches(3.5)), int(Inches(1.2)),
                 size=44, bold=True, color=C_BLUE, align="center")
            _txt(slide, kpi.get("label", ""), x, y + int(Inches(1.0)), int(Inches(3.5)), int(Inches(0.4)),
                 size=13, color=C_DGRAY, align="center")
    except Exception as e:
        print(f"[slide_kpis error] {e}")


def _slide_funnel(prs, layout):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "全渠道消費者旅程佈局")
        stages = [
            ("TOFU 認知", "解決「我不認識你」的問題", C_NAVY),
            ("MOFU 比較", "解決「為什麼選你」的問題", (0x15, 0x65, 0xC0)),
            ("BOFU 購買", "解決「現在就下單」的問題", (0x0D, 0x47, 0xA1)),
            ("Loyalty 回購", "解決「品牌一生相隨」的問題", C_BLUE),
        ]
        w = int(Inches(10))
        for i, (stage, desc, color) in enumerate(stages):
            wi = w - int(Inches(i * 0.5))
            xi = int(Inches(1.67)) + int(Inches(i * 0.25))
            yi = int(Inches(1.4)) + i * int(Inches(1.3))
            _rect(slide, xi, yi, wi, int(Inches(1.1)), fill=color)
            _txt(slide, stage, xi, yi + int(Inches(0.15)), wi, int(Inches(0.45)),
                 size=16, bold=True, color=C_WHITE, align="center")
            _txt(slide, desc, xi, yi + int(Inches(0.55)), wi, int(Inches(0.4)),
                 size=12, color=C_LBLUE, align="center")
    except Exception as e:
        print(f"[slide_funnel error] {e}")


def _slide_strategy(prs, layout, strategy):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "媒體策略戰術架構")
        headers = ["階段", "媒體工具", "溝通核心"]
        rows = [(s.get("stage", ""), s.get("tools", ""), s.get("message", "")) for s in strategy]
        _table(slide, headers, rows, int(Inches(0.5)), int(Inches(1.3)), int(Inches(12.33)), int(Inches(5.5)))
    except Exception as e:
        print(f"[slide_strategy error] {e}")


def _slide_budget(prs, layout, budget_rows, budget_total, budget_note):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, f"預算資源分配（月預算 {budget_total}萬）")
        bar_max_w = int(Inches(6.5))
        by = int(Inches(1.4))
        for channel, pct, amount in budget_rows:
            pct_num = int(re.sub(r"[^\d]", "", str(pct)) or "0")
            bar_w = int(bar_max_w * pct_num // 35)
            bar_w = max(bar_w, 1)
            _txt(slide, channel, int(Inches(0.5)), by, int(Inches(2.0)), int(Inches(0.38)), size=12, color=C_DGRAY)
            _rect(slide, int(Inches(2.6)), by + int(Inches(0.04)), min(bar_w, bar_max_w), int(Inches(0.3)), fill=C_NAVY)
            _txt(slide, f"{pct} ({amount})", int(Inches(2.7)), by + int(Inches(0.04)), int(Inches(3.0)), int(Inches(0.3)),
                 size=11, bold=True, color=C_WHITE)
            by += int(Inches(0.78))
        _rect(slide, int(Inches(9.8)), int(Inches(1.3)), int(Inches(3.1)), int(Inches(5.5)), fill=C_LGRAY, border=C_BORDER)
        _txt(slide, "策略說明：", int(Inches(10.0)), int(Inches(1.45)), int(Inches(2.7)), int(Inches(0.35)),
             size=13, bold=True, color=C_NAVY)
        notes = [s.strip() for s in budget_note.split("。") if s.strip()][:4]
        for ni, note in enumerate(notes):
            _txt(slide, f"{ni+1}. {note}。", int(Inches(10.0)), int(Inches(1.9)) + ni * int(Inches(1.0)),
                 int(Inches(2.7)), int(Inches(0.9)), size=11, color=C_DGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_budget error] {e}")


def _slide_meta(prs, layout, audiences):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "Meta 獲客主力策略")
        _txt(slide, "運用 Advantage+ 購物廣告 (ASC) 結合動態創意，自動優化轉化率。",
             int(Inches(0.5)), int(Inches(6.7)), int(Inches(12.33)), int(Inches(0.4)),
             size=11, italic=True, color=C_MGRAY, align="center")
        card_w = int(Inches(3.8))
        xs = [int(Inches(0.4)), int(Inches(4.77)), int(Inches(9.14))]
        for i, aud in enumerate(audiences[:3]):
            x = xs[i]
            _rect(slide, x, int(Inches(1.3)), card_w, int(Inches(5.2)), fill=C_WHITE, border=C_BORDER)
            _txt(slide, aud.get("name", ""), x, int(Inches(1.6)), card_w, int(Inches(0.5)),
                 size=16, bold=True, color=C_NAVY, align="center")
            _txt(slide, f"主打：{aud.get('products', '')}", x + int(Inches(0.2)), int(Inches(2.3)),
                 card_w - int(Inches(0.4)), int(Inches(0.4)), size=13, color=C_DGRAY)
            _txt(slide, f"創意：{aud.get('creative', '')}", x + int(Inches(0.2)), int(Inches(2.85)),
                 card_w - int(Inches(0.4)), int(Inches(0.5)), size=12, color=C_MGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_meta error] {e}")


def _slide_google(prs, layout, keywords):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "Google 高轉換關鍵字佈局")
        _rect(slide, int(Inches(0.5)), int(Inches(1.2)), int(Inches(12.33)), int(Inches(0.55)), fill=C_LGRAY, border=C_BORDER)
        _txt(slide, "核心邏輯：購買意圖強烈，必須佔據搜尋結果第一屏。",
             int(Inches(0.7)), int(Inches(1.28)), int(Inches(12.0)), int(Inches(0.38)),
             size=13, bold=True, color=C_NAVY)
        headers = ["關鍵字類別", "推薦關鍵字", "預期效果"]
        rows = [(k.get("type", ""), k.get("keywords", ""), k.get("effect", "")) for k in keywords]
        _table(slide, headers, rows, int(Inches(0.5)), int(Inches(2.0)), int(Inches(12.33)), int(Inches(4.7)))
    except Exception as e:
        print(f"[slide_google error] {e}")


def _slide_youtube(prs, layout, experts):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "YouTube 專家權威內容")
        _rect(slide, int(Inches(6.5)), int(Inches(1.25)), int(Inches(6.4)), int(Inches(5.5)), fill=C_LGRAY, border=C_BORDER)
        _txt(slide, "建立「不可忽視」的品牌信任", int(Inches(6.7)), int(Inches(1.55)), int(Inches(5.8)), int(Inches(0.5)),
             size=15, bold=True, color=C_NAVY)
        for i, exp in enumerate(experts[:3]):
            _txt(slide, f"· {exp.get('role', '')}：{exp.get('content', '')}", int(Inches(6.7)),
                 int(Inches(2.2)) + i * int(Inches(0.75)), int(Inches(5.8)), int(Inches(0.6)), size=13, color=C_DGRAY, wrap=True)
        _rect(slide, int(Inches(6.7)), int(Inches(5.55)), int(Inches(5.8)), int(Inches(0.75)), fill=C_NAVY)
        _txt(slide, "目標：將 YouTube 打造為品牌的「數位講堂」，提高信任轉化。",
             int(Inches(6.85)), int(Inches(5.65)), int(Inches(5.5)), int(Inches(0.55)), size=12, color=C_WHITE, wrap=True)
        _rect(slide, int(Inches(0.5)), int(Inches(1.55)), int(Inches(5.5)), int(Inches(4.5)), fill=C_LBLUE, border=C_BORDER)
        _txt(slide, "專家內容行銷", int(Inches(0.5)), int(Inches(5.8)), int(Inches(5.5)), int(Inches(0.4)),
             size=13, color=C_MGRAY, align="center")
        _txt(slide, "[YT]", int(Inches(1.5)), int(Inches(2.8)), int(Inches(3.5)), int(Inches(1.5)),
             size=48, align="center", color=C_NAVY)
    except Exception as e:
        print(f"[slide_youtube error] {e}")


def _slide_kol(prs, layout, tiers):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "KOL 三層式矩陣佈局")
        card_w = int(Inches(3.8))
        xs = [int(Inches(0.4)), int(Inches(4.77)), int(Inches(9.14))]
        bgs = [C_LGRAY, C_LBLUE, C_LGRAY]
        for i, tier in enumerate(tiers[:3]):
            x = xs[i]
            _rect(slide, x, int(Inches(1.3)), card_w, int(Inches(5.2)), fill=bgs[i], border=C_BORDER)
            _txt(slide, tier.get("tier", ""), x, int(Inches(1.7)), card_w, int(Inches(0.55)),
                 size=16, bold=True, color=C_NAVY, align="center")
            _txt(slide, f"目的：{tier.get('purpose', '')}", x, int(Inches(2.35)), card_w, int(Inches(0.4)),
                 size=13, bold=True, color=C_BLUE, align="center")
            _txt(slide, tier.get("desc", ""), x + int(Inches(0.2)), int(Inches(3.0)),
                 card_w - int(Inches(0.4)), int(Inches(2.0)), size=13, color=C_DGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_kol error] {e}")


def _slide_crm(prs, layout, steps):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "LINE 自動化 CRM 回購旅程")
        _rect(slide, int(Inches(0.5)), int(Inches(1.15)), int(Inches(0.08)), int(Inches(0.3)), fill=C_BLUE)
        _txt(slide, "保健品牌的勝負在「回購」", int(Inches(0.7)), int(Inches(1.15)), int(Inches(6.0)), int(Inches(0.35)),
             size=14, bold=True, color=C_BLUE)
        n = len(steps[:5])
        step_w = int(Inches(12.33 / n))
        _rect(slide, int(Inches(0.5)), int(Inches(3.55)), int(Inches(12.33)), int(Inches(0.06)), fill=C_BORDER)
        for i, step in enumerate(steps[:5]):
            x = int(Inches(0.5)) + i * step_w
            cx = x + step_w // 2 - int(Inches(0.35))
            _rect(slide, cx, int(Inches(3.1)), int(Inches(0.7)), int(Inches(0.7)), fill=C_NAVY)
            _txt(slide, step.get("day", ""), x, int(Inches(2.4)), step_w, int(Inches(0.35)),
                 size=11, bold=True, color=C_NAVY, align="center")
            _txt(slide, step.get("title", ""), x, int(Inches(4.0)), step_w, int(Inches(0.45)),
                 size=13, bold=True, color=C_DGRAY, align="center")
            _txt(slide, step.get("desc", ""), x, int(Inches(4.5)), step_w, int(Inches(0.45)),
                 size=11, color=C_MGRAY, align="center")
        _rect(slide, int(Inches(0.5)), int(Inches(5.3)), int(Inches(12.33)), int(Inches(0.85)), fill=C_LGRAY, border=C_BORDER)
        _txt(slide, "重點：透過分眾標籤（已購產品、性別、年齡）實現精準自動化推播，降低手動操作成本。",
             int(Inches(0.7)), int(Inches(5.4)), int(Inches(12.0)), int(Inches(0.6)), size=12, bold=True, color=C_NAVY, wrap=True)
    except Exception as e:
        print(f"[slide_crm error] {e}")


def _slide_resources(prs, layout, must_buy, bonus):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "必備媒體資源建議")
        _txt(slide, "主流必買工具 (Must Buy)", int(Inches(0.5)), int(Inches(1.2)), int(Inches(5.8)), int(Inches(0.45)),
             size=15, bold=True, color=C_NAVY)
        _rect(slide, int(Inches(0.5)), int(Inches(1.65)), int(Inches(5.8)), int(Inches(0.05)), fill=C_NAVY)
        _bullets(slide, must_buy, int(Inches(0.5)), int(Inches(1.85)), int(Inches(5.8)), int(Inches(4.5)), size=13)
        _txt(slide, "加分資源 (Bonus)", int(Inches(7.1)), int(Inches(1.2)), int(Inches(5.8)), int(Inches(0.45)),
             size=15, bold=True, color=C_NAVY)
        _rect(slide, int(Inches(7.1)), int(Inches(1.65)), int(Inches(5.8)), int(Inches(0.05)), fill=C_BLUE)
        _bullets(slide, bonus, int(Inches(7.1)), int(Inches(1.85)), int(Inches(5.8)), int(Inches(4.5)), size=13)
    except Exception as e:
        print(f"[slide_resources error] {e}")


def _slide_quarterly(prs, layout, plan):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "年度實施藍圖")
        headers = ["季度", "核心目標", "配置策略"]
        rows = [(p.get("quarter", ""), p.get("goal", ""), p.get("strategy", "")) for p in plan]
        _table(slide, headers, rows, int(Inches(0.5)), int(Inches(1.3)), int(Inches(12.33)), int(Inches(5.5)))
    except Exception as e:
        print(f"[slide_quarterly error] {e}")


def _slide_closing(prs, layout, client_name, closing_msg, year):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _set_bg(slide, C_NAVY)
        _txt(slide, f"{client_name} Growth Blueprint {year}", int(Inches(1)), int(Inches(1.0)), int(Inches(11.33)), int(Inches(1.2)),
             size=32, bold=True, color=C_WHITE, align="center")
        _rect(slide, int(Inches(4)), int(Inches(2.3)), int(Inches(5.33)), int(Inches(0.07)), fill=C_BLUE)
        _txt(slide, closing_msg, int(Inches(1.5)), int(Inches(2.7)), int(Inches(10.33)), int(Inches(2.0)),
             size=16, color=C_WHITE, align="center", wrap=True)
        _txt(slide, "「內容 x 搜尋 x CRM x 會員經營」全漏斗成長模式",
             int(Inches(1.5)), int(Inches(5.0)), int(Inches(10.33)), int(Inches(0.7)),
             size=18, bold=True, color=C_BLUE, align="center")
        _txt(slide, "潮網科技 Wavenet Technology", int(Inches(3)), int(Inches(6.5)), int(Inches(7.33)), int(Inches(0.4)),
             size=12, color=C_MGRAY, align="center")
    except Exception as e:
        print(f"[slide_closing error] {e}")


# ── B2B Biotech Individual slides ─────────────────────────────────────────────

def _slide_b2b_cover(prs, layout, client_name, subtitle, year):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _set_bg(slide, C_NAVY)
        _rect(slide, int(Inches(0)), int(Inches(3.2)), int(Inches(13.33)), int(Inches(0.08)), fill=C_BLUE)
        _txt(slide, f"{client_name} {year}", int(Inches(1.5)), int(Inches(1.0)), int(Inches(10)), int(Inches(1.2)),
             size=40, bold=True, color=C_WHITE, align="center")
        _txt(slide, "B2B 生技產業媒體提案", int(Inches(1.5)), int(Inches(2.1)), int(Inches(10)), int(Inches(0.9)),
             size=28, bold=True, color=C_BLUE, align="center")
        _txt(slide, subtitle, int(Inches(2)), int(Inches(3.45)), int(Inches(9.33)), int(Inches(0.7)),
             size=16, color=C_WHITE, align="center", italic=True)
        _txt(slide, "Powered by 潮網科技 Wavenet Technology", int(Inches(3)), int(Inches(6.8)), int(Inches(7.33)), int(Inches(0.4)),
             size=11, color=C_MGRAY, align="center")
    except Exception as e:
        print(f"[slide_b2b_cover error] {e}")


def _slide_b2b_tech(prs, layout, tech_highlights, business_model):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "技術與品牌現況分析")
        _rect(slide, int(Inches(0.4)), int(Inches(1.25)), int(Inches(5.8)), int(Inches(5.5)), fill=C_LGRAY, border=C_BORDER)
        _bullets(slide, tech_highlights, int(Inches(0.6)), int(Inches(1.45)), int(Inches(5.4)), int(Inches(5.0)),
                 title="核心技術亮點", title_size=16, size=14)
        _rect(slide, int(Inches(7.1)), int(Inches(1.25)), int(Inches(5.8)), int(Inches(5.5)), fill=C_LBLUE, border=C_BORDER)
        _txt(slide, "商業模式", int(Inches(7.3)), int(Inches(1.5)), int(Inches(5.4)), int(Inches(0.5)),
             size=16, bold=True, color=C_NAVY)
        _txt(slide, business_model, int(Inches(7.3)), int(Inches(2.1)), int(Inches(5.4)), int(Inches(2.5)),
             size=13, color=C_DGRAY, wrap=True)
        _rect(slide, int(Inches(7.3)), int(Inches(5.0)), int(Inches(5.4)), int(Inches(0.6)), fill=C_LBLUE2)
        _txt(slide, "核心競爭力：專利技術護城河、臨床數據支撐。", int(Inches(7.5)), int(Inches(5.05)), int(Inches(5.0)), int(Inches(0.5)),
             size=12, bold=True, color=C_NAVY)
    except Exception as e:
        print(f"[slide_b2b_tech error] {e}")


def _slide_b2b_challenges(prs, layout, challenges):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "挑戰診斷")
        card_w = int(Inches(3.8))
        xs = [int(Inches(0.4)), int(Inches(4.77)), int(Inches(9.14))]
        for i, ch in enumerate(challenges[:3]):
            x = xs[i]
            _rect(slide, x, int(Inches(1.3)), card_w, int(Inches(5.2)), fill=C_LORANGE, border=C_BORDER)
            _txt(slide, ch.get("title", ""), x, int(Inches(1.7)), card_w, int(Inches(0.55)),
                 size=16, bold=True, color=C_ORANGE, align="center")
            _txt(slide, ch.get("desc", ""), x + int(Inches(0.2)), int(Inches(2.5)),
                 card_w - int(Inches(0.4)), int(Inches(3.5)), size=13, color=C_DGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_b2b_challenges error] {e}")


def _slide_b2b_journey(prs, layout, journey):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "B2B 決策者旅程")
        stages = journey[:4]
        colors = [C_NAVY, (0x15, 0x65, 0xC0), (0x0D, 0x47, 0xA1), C_BLUE]
        step_w = int(Inches(12.33 / max(len(stages), 1)))
        # Timeline
        _rect(slide, int(Inches(0.5)), int(Inches(3.6)), int(Inches(12.33)), int(Inches(0.08)), fill=C_BORDER)
        for i, jrn in enumerate(stages):
            x = int(Inches(0.5)) + i * step_w
            cx = x + step_w // 2 - int(Inches(0.4))
            color = colors[i % len(colors)]
            # Stage box top
            _rect(slide, x + int(Inches(0.1)), int(Inches(1.3)), step_w - int(Inches(0.2)), int(Inches(1.6)), fill=color)
            _txt(slide, jrn.get("stage", ""), x, int(Inches(1.4)), step_w, int(Inches(0.45)),
                 size=14, bold=True, color=C_WHITE, align="center")
            _txt(slide, jrn.get("action", ""), x, int(Inches(1.9)), step_w, int(Inches(0.45)),
                 size=12, color=C_LBLUE, align="center")
            # Circle on timeline
            _rect(slide, cx, int(Inches(3.3)), int(Inches(0.8)), int(Inches(0.8)), fill=color)
            # Description below
            _txt(slide, jrn.get("desc", ""), x + int(Inches(0.1)), int(Inches(4.2)),
                 step_w - int(Inches(0.2)), int(Inches(2.0)), size=12, color=C_DGRAY, align="center", wrap=True)
    except Exception as e:
        print(f"[slide_b2b_journey error] {e}")


def _slide_b2b_budget(prs, layout, budget_total):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, f"預算資源分配（年度預算 {budget_total}萬）")
        budget_rows = [
            ("LinkedIn ABM廣告",   "35%", f"{int(int(budget_total)*0.35):,}"),
            ("Google Search/SEO", "25%", f"{int(int(budget_total)*0.25):,}"),
            ("生技媒體公關",        "15%", f"{int(int(budget_total)*0.15):,}"),
            ("會展行銷",            "15%", f"{int(int(budget_total)*0.15):,}"),
            ("內容製作",            "10%", f"{int(int(budget_total)*0.10):,}"),
        ]
        bar_max_w = int(Inches(6.5))
        by = int(Inches(1.4))
        for channel, pct, amount in budget_rows:
            pct_num = int(re.sub(r"[^\d]", "", str(pct)) or "0")
            bar_w = int(bar_max_w * pct_num // 35)
            bar_w = max(bar_w, 1)
            _txt(slide, channel, int(Inches(0.5)), by, int(Inches(2.2)), int(Inches(0.38)), size=12, color=C_DGRAY)
            _rect(slide, int(Inches(2.8)), by + int(Inches(0.04)), min(bar_w, bar_max_w), int(Inches(0.3)), fill=C_NAVY)
            _txt(slide, f"{pct} ({amount}萬)", int(Inches(2.9)), by + int(Inches(0.04)), int(Inches(3.0)), int(Inches(0.3)),
                 size=11, bold=True, color=C_WHITE)
            by += int(Inches(0.78))
        _rect(slide, int(Inches(9.8)), int(Inches(1.3)), int(Inches(3.1)), int(Inches(5.5)), fill=C_LGRAY, border=C_BORDER)
        _txt(slide, "配置原則：", int(Inches(10.0)), int(Inches(1.45)), int(Inches(2.7)), int(Inches(0.35)),
             size=13, bold=True, color=C_NAVY)
        notes = ["LinkedIn 為 B2B 核心觸及管道", "SEO 建立長期專業形象", "展會是關係建立關鍵", "內容資產持續累積信任"]
        for ni, note in enumerate(notes):
            _txt(slide, f"{ni+1}. {note}", int(Inches(10.0)), int(Inches(1.9)) + ni * int(Inches(1.0)),
                 int(Inches(2.7)), int(Inches(0.9)), size=11, color=C_DGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_b2b_budget error] {e}")


def _slide_b2b_linkedin(prs, layout, linkedin):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "LinkedIn ABM 精準策略")
        sections = [
            ("目標職稱", linkedin.get("roles", []), C_LBLUE),
            ("目標公司類型", linkedin.get("companies", []), C_LGRAY),
            ("廣告形式", linkedin.get("ad_formats", []), C_LBLUE),
        ]
        card_w = int(Inches(3.8))
        xs = [int(Inches(0.4)), int(Inches(4.77)), int(Inches(9.14))]
        for i, (title, items, bg) in enumerate(sections):
            x = xs[i]
            _rect(slide, x, int(Inches(1.3)), card_w, int(Inches(5.2)), fill=bg, border=C_BORDER)
            _txt(slide, title, x, int(Inches(1.5)), card_w, int(Inches(0.5)),
                 size=15, bold=True, color=C_NAVY, align="center")
            _rect(slide, x + int(Inches(0.3)), int(Inches(2.0)), card_w - int(Inches(0.6)), int(Inches(0.05)), fill=C_NAVY)
            _bullets(slide, items, x + int(Inches(0.2)), int(Inches(2.2)), card_w - int(Inches(0.4)), int(Inches(3.8)), size=13)
    except Exception as e:
        print(f"[slide_b2b_linkedin error] {e}")


def _slide_b2b_seo(prs, layout, seo_keywords):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "Google & SEO 策略")
        _rect(slide, int(Inches(0.5)), int(Inches(1.2)), int(Inches(12.33)), int(Inches(0.55)), fill=C_LGRAY, border=C_BORDER)
        _txt(slide, "核心邏輯：生技決策者在評估供應商時高度依賴搜尋，必須佔據專業關鍵字首頁。",
             int(Inches(0.7)), int(Inches(1.28)), int(Inches(12.0)), int(Inches(0.38)),
             size=13, bold=True, color=C_NAVY)
        headers = ["關鍵字類別", "推薦關鍵字", "搜尋意圖"]
        rows = [(k.get("type", ""), k.get("keywords", ""), k.get("intent", "")) for k in seo_keywords]
        _table(slide, headers, rows, int(Inches(0.5)), int(Inches(2.0)), int(Inches(12.33)), int(Inches(4.7)))
    except Exception as e:
        print(f"[slide_b2b_seo error] {e}")


def _slide_b2b_thought_leadership(prs, layout, thought_leadership):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "思想領導力內容策略")
        card_w = int(Inches(3.8))
        xs = [int(Inches(0.4)), int(Inches(4.77)), int(Inches(9.14))]
        bgs = [C_LGRAY, C_LBLUE, C_LGRAY]
        for i, tl in enumerate(thought_leadership[:3]):
            x = xs[i]
            _rect(slide, x, int(Inches(1.3)), card_w, int(Inches(5.2)), fill=bgs[i], border=C_BORDER)
            _txt(slide, tl.get("format", ""), x, int(Inches(1.55)), card_w, int(Inches(0.55)),
                 size=14, bold=True, color=C_NAVY, align="center")
            _rect(slide, x + int(Inches(0.3)), int(Inches(2.1)), card_w - int(Inches(0.6)), int(Inches(0.05)), fill=C_NAVY)
            _txt(slide, f"主題：{tl.get('topic', '')}", x + int(Inches(0.2)), int(Inches(2.3)),
                 card_w - int(Inches(0.4)), int(Inches(1.5)), size=13, color=C_DGRAY, wrap=True)
            _txt(slide, f"目標：{tl.get('goal', '')}", x + int(Inches(0.2)), int(Inches(4.0)),
                 card_w - int(Inches(0.4)), int(Inches(1.5)), size=12, bold=True, color=C_BLUE, wrap=True)
    except Exception as e:
        print(f"[slide_b2b_thought_leadership error] {e}")


def _slide_b2b_events(prs, layout, events):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "會展行銷策略")
        for i, evt in enumerate(events[:4]):
            y = int(Inches(1.35)) + i * int(Inches(1.45))
            _rect(slide, int(Inches(0.5)), y, int(Inches(12.33)), int(Inches(1.25)), fill=C_LGRAY, border=C_BORDER)
            _txt(slide, evt.get("name", ""), int(Inches(0.7)), y + int(Inches(0.1)), int(Inches(3.5)), int(Inches(0.4)),
                 size=14, bold=True, color=C_NAVY)
            _txt(slide, f"策略：{evt.get('strategy', '')}", int(Inches(4.3)), y + int(Inches(0.1)), int(Inches(4.0)), int(Inches(0.4)),
                 size=13, color=C_DGRAY)
            _txt(slide, f"戰術：{evt.get('tactic', '')}", int(Inches(8.5)), y + int(Inches(0.1)), int(Inches(4.0)), int(Inches(0.55)),
                 size=12, color=C_MGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_b2b_events error] {e}")


def _slide_b2b_content_assets(prs, layout, content_assets):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _header(slide, "內容資產規劃")
        _txt(slide, "建立完整的數位內容資產庫，支撐各階段買家旅程的溝通需求。",
             int(Inches(0.5)), int(Inches(1.15)), int(Inches(12.33)), int(Inches(0.4)),
             size=13, italic=True, color=C_MGRAY)
        for i, asset in enumerate(content_assets[:5]):
            y = int(Inches(1.75)) + i * int(Inches(1.0))
            _rect(slide, int(Inches(0.5)), y, int(Inches(12.33)), int(Inches(0.85)), fill=C_LBLUE if i % 2 == 0 else C_LGRAY, border=C_BORDER)
            _txt(slide, asset.get("type", ""), int(Inches(0.7)), y + int(Inches(0.12)), int(Inches(2.5)), int(Inches(0.5)),
                 size=14, bold=True, color=C_NAVY)
            _txt(slide, asset.get("desc", ""), int(Inches(3.4)), y + int(Inches(0.12)), int(Inches(9.0)), int(Inches(0.55)),
                 size=13, color=C_DGRAY, wrap=True)
    except Exception as e:
        print(f"[slide_b2b_content_assets error] {e}")


def _slide_b2b_closing(prs, layout, client_name, closing_msg, year):
    try:
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        _set_bg(slide, C_NAVY)
        _txt(slide, f"{client_name} B2B Growth Blueprint {year}", int(Inches(1)), int(Inches(1.0)), int(Inches(11.33)), int(Inches(1.2)),
             size=28, bold=True, color=C_WHITE, align="center")
        _rect(slide, int(Inches(4)), int(Inches(2.3)), int(Inches(5.33)), int(Inches(0.07)), fill=C_BLUE)
        _txt(slide, closing_msg, int(Inches(1.5)), int(Inches(2.7)), int(Inches(10.33)), int(Inches(2.0)),
             size=16, color=C_WHITE, align="center", wrap=True)
        _txt(slide, "「LinkedIn x SEO x 會展 x 思想領導力」B2B 全漏斗成長模式",
             int(Inches(1.5)), int(Inches(5.0)), int(Inches(10.33)), int(Inches(0.7)),
             size=16, bold=True, color=C_BLUE, align="center")
        _txt(slide, "潮網科技 Wavenet Technology", int(Inches(3)), int(Inches(6.5)), int(Inches(7.33)), int(Inches(0.4)),
             size=12, color=C_MGRAY, align="center")
    except Exception as e:
        print(f"[slide_b2b_closing error] {e}")


# ── Build PPTX ────────────────────────────────────────────────────────────────

def _build_pptx(content: dict, body: ProposalRequest) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Emu

    budget_num = int(re.sub(r"[^\d]", "", str(body.monthly_budget)) or "100")

    prs = Presentation()
    prs.slide_width  = int(Inches(13.33))
    prs.slide_height = int(Inches(7.5))
    layout = prs.slide_layouts[6]  # blank

    _build_pptx_unified(prs, layout, content, body, budget_num)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_pptx_b2c(prs, layout, content: dict, body: ProposalRequest, budget_num: int):
    year = body.year
    client_name = body.client_name

    # Budget rows for budget slide
    budget_rows = [
        ("Meta FB/IG 廣告",    "35%", f"{int(budget_num*0.35)}"),
        ("Google Ads / PMax",  "25%", f"{int(budget_num*0.25)}"),
        ("YouTube 影音",        "15%", f"{int(budget_num*0.15)}"),
        ("LINE CRM / LAP",     "15%", f"{int(budget_num*0.15)}"),
        ("KOL / 內容行銷",     "10%", f"{int(budget_num*0.10)}"),
    ]
    budget_note = (
        f"Meta 廣告佔最大比重，以 ASC 自動優化為主。"
        f"Google PMax 延伸流量。"
        f"LINE 強化回購旅程。"
        f"KOL 製造社群口碑。"
    )

    # Slide 1: Cover
    _slide_cover(prs, layout, client_name, content.get("subtitle", ""), year)

    # Slide 2: Brand Analysis
    _slide_brand(prs, layout,
                 content.get("brand_strengths", []),
                 content.get("brand_d2c", ""))

    # Slide 3: Market Segments
    _slide_market(prs, layout, content.get("market_segments", []), year)

    # Slide 4: Problems
    _slide_problems(prs, layout, content.get("problems", []))

    # Slide 5: KPIs
    _slide_kpis(prs, layout, content.get("kpis", []), year)

    # Slide 6: Funnel
    _slide_funnel(prs, layout)

    # Slide 7: Strategy Table
    _slide_strategy(prs, layout, content.get("media_strategy", []))

    # Slide 8: Budget
    _slide_budget(prs, layout, budget_rows, str(budget_num), budget_note)

    # Slide 9: Meta
    _slide_meta(prs, layout, content.get("meta_audiences", []))

    # Slide 10: Google
    _slide_google(prs, layout, content.get("google_keywords", []))

    # Slide 11: YouTube
    _slide_youtube(prs, layout, content.get("youtube_experts", []))

    # Slide 12: KOL
    _slide_kol(prs, layout, content.get("kol_tiers", []))

    # Slide 13: CRM
    _slide_crm(prs, layout, content.get("crm_steps", []))

    # Slide 14: Resources
    _slide_resources(prs, layout,
                     content.get("must_buy", []),
                     content.get("bonus_resources", []))

    # Slide 15: Quarterly Plan
    _slide_quarterly(prs, layout, content.get("quarterly_plan", []))

    # Slide 16: Closing
    _slide_closing(prs, layout, client_name, content.get("closing_message", ""), year)


def _build_pptx_b2b_biotech(prs, layout, content: dict, body: ProposalRequest, budget_num: int):
    year = body.year
    client_name = body.client_name

    # Slide 1: Cover
    _slide_b2b_cover(prs, layout, client_name, content.get("subtitle", ""), year)

    # Slide 2: Tech & Brand Analysis
    _slide_b2b_tech(prs, layout,
                    content.get("tech_highlights", []),
                    content.get("business_model", ""))

    # Slide 3: Market Problems
    _slide_problems(prs, layout, content.get("market_problems", []), title="市場痛點分析")

    # Slide 4: Challenges Diagnosis
    _slide_b2b_challenges(prs, layout, content.get("challenges", []))

    # Slide 5: Strategic KPIs
    _slide_kpis(prs, layout, content.get("kpis", []), year)

    # Slide 6: B2B Decision Journey
    _slide_b2b_journey(prs, layout, content.get("b2b_journey", []))

    # Slide 7: Media Strategy Table
    _slide_strategy(prs, layout, content.get("media_strategy", []))

    # Slide 8: Budget Allocation
    _slide_b2b_budget(prs, layout, str(budget_num))

    # Slide 9: LinkedIn ABM Strategy
    _slide_b2b_linkedin(prs, layout, content.get("linkedin_targeting", {}))

    # Slide 10: Google & SEO Strategy
    _slide_b2b_seo(prs, layout, content.get("seo_keywords", []))

    # Slide 11: Thought Leadership
    _slide_b2b_thought_leadership(prs, layout, content.get("thought_leadership", []))

    # Slide 12: Trade Show Strategy
    _slide_b2b_events(prs, layout, content.get("events", []))

    # Slide 13: Content Assets
    _slide_b2b_content_assets(prs, layout, content.get("content_assets", []))

    # Slide 14: Media Resources
    _slide_resources(prs, layout,
                     content.get("must_buy", []),
                     content.get("bonus_resources", []))

    # Slide 15: Quarterly Roadmap
    _slide_quarterly(prs, layout, content.get("quarterly_plan", []))

    # Slide 16: Closing
    _slide_b2b_closing(prs, layout, client_name, content.get("closing_message", ""), year)


def _build_pptx_unified(prs, layout, content: dict, body: ProposalRequest, budget_num: int):
    year = body.year
    client_name = body.client_name

    budget_rows = [
        ("Meta FB/IG 廣告",   "30%", f"{int(budget_num*0.30)}"),
        ("Google Ads / PMax", "20%", f"{int(budget_num*0.20)}"),
        ("YouTube 影音",       "10%", f"{int(budget_num*0.10)}"),
        ("LinkedIn 廣告",     "15%", f"{int(budget_num*0.15)}"),
        ("LINE CRM / LAP",    "10%", f"{int(budget_num*0.10)}"),
        ("KOL / 內容行銷",    "10%", f"{int(budget_num*0.10)}"),
        ("展會 / 研討會",      "5%",  f"{int(budget_num*0.05)}"),
    ]
    budget_note = (
        "Meta + Google 主力效果媒體，覆蓋消費者全旅程。"
        "LinkedIn 精準觸及 B2B 決策者。"
        "KOL + 內容行銷建立品牌信任。"
        "LINE CRM 強化回購與客戶留存。"
    )

    _slide_cover(prs, layout, client_name, content.get("subtitle", ""), year)
    _slide_brand(prs, layout, content.get("brand_strengths", []), content.get("brand_d2c", ""))
    _slide_market(prs, layout, content.get("market_segments", []), year)
    _slide_problems(prs, layout, content.get("problems", []))
    _slide_kpis(prs, layout, content.get("kpis", []), year)
    _slide_funnel(prs, layout)
    _slide_strategy(prs, layout, content.get("media_strategy", []))
    _slide_budget(prs, layout, budget_rows, str(budget_num), budget_note)
    _slide_meta(prs, layout, content.get("meta_audiences", []))
    _slide_google(prs, layout, content.get("google_keywords", []))
    _slide_youtube(prs, layout, content.get("youtube_experts", []))
    _slide_kol(prs, layout, content.get("kol_tiers", []))
    _slide_b2b_linkedin(prs, layout, content.get("linkedin_targeting", {}))
    _slide_b2b_seo(prs, layout, content.get("seo_keywords", []))
    _slide_b2b_thought_leadership(prs, layout, content.get("thought_leadership", []))
    _slide_b2b_events(prs, layout, content.get("events", []))
    _slide_crm(prs, layout, content.get("crm_steps", []))
    _slide_b2b_content_assets(prs, layout, content.get("content_assets", []))
    _slide_resources(prs, layout, content.get("must_buy", []), content.get("bonus_resources", []))
    _slide_quarterly(prs, layout, content.get("quarterly_plan", []))
    _slide_closing(prs, layout, client_name, content.get("closing_message", ""), year)
